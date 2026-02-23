"""
SaaS市場調査 コンサルスライド生成 (BCGスタイル)
python build_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ===== カラーパレット (BCG風) =====
NAVY   = RGBColor(0x00, 0x2B, 0x5C)   # 濃紺
BLUE   = RGBColor(0x00, 0x5A, 0xA7)   # 中青
LBLUE  = RGBColor(0xD6, 0xE8, 0xF5)   # 薄青（背景帯）
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)   # アクセント

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_rect(slide, x, y, w, h, fill_color, border=False, border_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border and border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=11, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True, font_name="Yu Gothic"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def header_bar(slide, title_text, subtitle_text=""):
    """上部ネイビーヘッダー帯"""
    bar_h = Inches(1.15)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, NAVY)
    # タイトル (24pt 統一 per MEMORY.md)
    add_text(slide, title_text,
             Inches(0.4), Inches(0.12), Inches(11.0), Inches(0.7),
             font_size=24, bold=True, color=WHITE)
    if subtitle_text:
        add_text(slide, subtitle_text,
                 Inches(0.4), Inches(0.78), Inches(11.0), Inches(0.3),
                 font_size=10, bold=False, color=LBLUE)

def footer(slide, source_text="出典: Gartner調査 / 各社公開情報"):
    """下部フッター"""
    y = SLIDE_H - Inches(0.3)
    add_rect(slide, 0, y, SLIDE_W, Inches(0.28), LGRAY)
    add_text(slide, source_text,
             Inches(0.3), y + Pt(3), Inches(10), Inches(0.25),
             font_size=7, color=GRAY)
    add_text(slide, "Confidential",
             Inches(11.5), y + Pt(3), Inches(1.5), Inches(0.25),
             font_size=7, color=GRAY, align=PP_ALIGN.RIGHT)

def section_label(slide, text, x, y, w=Inches(2.5)):
    """セクションラベル（小さい青帯）"""
    add_rect(slide, x, y, w, Inches(0.26), BLUE)
    add_text(slide, text, x + Inches(0.08), y + Pt(2),
             w - Inches(0.1), Inches(0.24),
             font_size=9, bold=True, color=WHITE)

# ============================================================
# スライド1: 表紙
# ============================================================
def slide_cover(prs):
    sl = blank_slide(prs)
    # 背景
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # 右側アクセント帯
    add_rect(sl, SLIDE_W - Inches(3.5), 0, Inches(3.5), SLIDE_H, BLUE)
    # タイトル
    add_text(sl, "日本のSaaS市場\n2025年 市場調査レポート",
             Inches(0.6), Inches(1.8), Inches(8.5), Inches(2.2),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # サブタイトル
    add_text(sl, "市場規模 / 主要プレイヤー / トレンド / ビジネス示唆",
             Inches(0.6), Inches(4.0), Inches(8.5), Inches(0.5),
             font_size=14, color=LBLUE, align=PP_ALIGN.LEFT)
    # 日付
    add_text(sl, "2026年2月",
             Inches(0.6), Inches(4.7), Inches(4.0), Inches(0.4),
             font_size=11, color=WHITE)
    # 右帯内テキスト
    add_text(sl, "Market\nResearch",
             SLIDE_W - Inches(3.1), Inches(3.0), Inches(2.8), Inches(1.5),
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# スライド2: エグゼクティブサマリー
# ============================================================
def slide_exec_summary(prs):
    sl = blank_slide(prs)
    header_bar(sl, "エグゼクティブサマリー",
               "Key Findings — 日本のSaaS市場 2025年")

    # 3ボックスレイアウト
    boxes = [
        ("📈 市場規模",
         "世界SaaS市場は2020年の1,030億ドルから2022年には1,450億ドルへ拡大。CAGR約11%で成長継続。日本市場もDX推進を背景に高成長フェーズにある。"),
        ("🏢 主要プレイヤー",
         "Tier1はSalesforce（CRM）・Microsoft（Power Platform）。Tier2はGoogle Workspace・Adobe Experience Cloud。国内ではマネーフォワード・freee・Sansanが台頭。"),
        ("🎯 ビジネス示唆",
         "AI統合SaaSへのシフトが加速。SMB向け低価格帯と大企業向けエンタープライズで二極化。垂直特化型（業界特化SaaS）に参入余地あり。"),
    ]

    box_w = Inches(3.8)
    box_h = Inches(4.5)
    gap   = Inches(0.35)
    start_x = Inches(0.35)
    start_y = Inches(1.35)

    for i, (ttl, body) in enumerate(boxes):
        bx = start_x + i * (box_w + gap)
        # 外枠
        add_rect(sl, bx, start_y, box_w, box_h, WHITE,
                 border=True, border_color=BLUE)
        # 上帯
        add_rect(sl, bx, start_y, box_w, Inches(0.45), BLUE)
        add_text(sl, ttl, bx + Inches(0.1), start_y + Pt(5),
                 box_w - Inches(0.2), Inches(0.38),
                 font_size=12, bold=True, color=WHITE)
        # 本文
        add_text(sl, body, bx + Inches(0.15), start_y + Inches(0.55),
                 box_w - Inches(0.3), box_h - Inches(0.7),
                 font_size=10.5, color=BLACK, wrap=True)

    footer(sl)

# ============================================================
# スライド3: 市場規模
# ============================================================
def slide_market_size(prs):
    sl = blank_slide(prs)
    header_bar(sl, "市場規模 — 世界SaaS市場は年率11%で成長中",
               "Market Size (Source: Gartner, 2022)")

    # 左: 数字ハイライト
    section_label(sl, "市場規模推移", Inches(0.4), Inches(1.35))

    data = [
        ("2020年", "1,030", "億ドル"),
        ("2022年", "1,450", "億ドル"),
        ("2025年(予測)", "~2,100", "億ドル"),
    ]
    bar_colors = [LBLUE, BLUE, NAVY]
    bar_x = Inches(0.4)
    bar_start_y = Inches(1.75)
    max_val = 2100

    for i, (yr, val, unit) in enumerate(data):
        by = bar_start_y + i * Inches(1.45)
        num = int(val.replace(",","").replace("~",""))
        bw  = Inches(7.0) * num / max_val
        add_rect(sl, bar_x, by + Inches(0.38), bw, Inches(0.55), bar_colors[i])
        add_text(sl, yr,    bar_x, by, Inches(2.0), Inches(0.35), font_size=10, bold=True)
        add_text(sl, f"{val} {unit}", bar_x + bw + Inches(0.1), by + Inches(0.38),
                 Inches(2.0), Inches(0.55), font_size=14, bold=True, color=NAVY)

    # 右: CAGR ボックス
    cx = Inches(9.0)
    cy = Inches(1.35)
    add_rect(sl, cx, cy, Inches(3.9), Inches(2.2), LBLUE)
    add_text(sl, "CAGR (2020→2022)", cx + Inches(0.15), cy + Inches(0.1),
             Inches(3.5), Inches(0.4), font_size=10, bold=True, color=NAVY)
    add_text(sl, "≈ 11%",   cx + Inches(0.3),  cy + Inches(0.5),
             Inches(3.0), Inches(0.9), font_size=40, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(sl, "年平均成長率",  cx + Inches(0.3),  cy + Inches(1.5),
             Inches(3.0), Inches(0.4), font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # 右: 日本市場メモ
    add_rect(sl, cx, cy + Inches(2.4), Inches(3.9), Inches(2.0), LGRAY)
    add_text(sl, "🇯🇵 日本市場",
             cx + Inches(0.15), cy + Inches(2.5), Inches(3.5), Inches(0.35),
             font_size=10, bold=True, color=NAVY)
    add_text(sl,
             "• DX推進法・政府クラウド移行方針を背景に急拡大\n"
             "• 国内SaaS市場は2024年に約1兆円超と推定\n"
             "• 中小企業への浸透が今後の成長ドライバー",
             cx + Inches(0.15), cy + Inches(2.9), Inches(3.6), Inches(1.3),
             font_size=9, color=BLACK)

    footer(sl)

# ============================================================
# スライド4: 主要プレイヤー
# ============================================================
def slide_players(prs):
    sl = blank_slide(prs)
    header_bar(sl, "主要プレイヤー — グローバル大手が市場を牽引",
               "Key Players (Tier Classification)")

    players = {
        "Tier 1\nグローバルリーダー": [
            ("Salesforce", "CRM・営業支援のデファクトスタンダード。国内大企業での導入率が高い。"),
            ("Microsoft",  "Microsoft 365・Teams・Power Platformで業務全般をカバー。"),
        ],
        "Tier 2\nメジャープレイヤー": [
            ("Google",  "Google Workspace（旧G Suite）でSMB市場を中心に拡大。"),
            ("Adobe",   "Experience Cloudでマーケティング・デジタル体験領域に特化。"),
        ],
        "Tier 3\n国内注目プレイヤー": [
            ("マネーフォワード", "クラウド会計・給与で中小企業向けに急成長。"),
            ("Sansan",          "名刺管理→営業DXへ領域拡張。エンタープライズ向け。"),
        ],
    }

    tier_colors = [NAVY, BLUE, RGBColor(0x4C, 0x8C, 0xC4)]
    col_w = Inches(4.1)
    col_h = Inches(5.2)
    gap   = Inches(0.2)
    sx    = Inches(0.25)
    sy    = Inches(1.3)

    for ci, (tier_name, items) in enumerate(players.items()):
        cx = sx + ci * (col_w + gap)
        # 列ヘッダー
        add_rect(sl, cx, sy, col_w, Inches(0.55), tier_colors[ci])
        add_text(sl, tier_name, cx + Inches(0.1), sy + Pt(4),
                 col_w - Inches(0.2), Inches(0.48),
                 font_size=10, bold=True, color=WHITE)
        # 各社カード
        for ri, (name, desc) in enumerate(items):
            ry = sy + Inches(0.65) + ri * Inches(2.2)
            add_rect(sl, cx, ry, col_w, Inches(2.1), WHITE,
                     border=True, border_color=tier_colors[ci])
            add_text(sl, name, cx + Inches(0.1), ry + Inches(0.08),
                     col_w - Inches(0.2), Inches(0.45),
                     font_size=13, bold=True, color=tier_colors[ci])
            add_text(sl, desc, cx + Inches(0.1), ry + Inches(0.55),
                     col_w - Inches(0.2), Inches(1.4),
                     font_size=9.5, color=BLACK)

    footer(sl)

# ============================================================
# スライド5: トレンド
# ============================================================
def slide_trends(prs):
    sl = blank_slide(prs)
    header_bar(sl, "主要トレンド — AI統合・業界特化・コスト最適化が加速",
               "Key Trends 2025")

    trends = [
        ("① AI統合SaaSの台頭",
         "ChatGPT等のLLMがSaaSに組み込まれ、業務自動化が急加速。Salesforce Einstein・Microsoft Copilotが先行。",
         "High"),
        ("② 垂直特化（バーティカルSaaS）",
         "医療・建設・製造など業界特化型SaaSが急成長。汎用SaaSとの差別化で高い解約率低下を実現。",
         "High"),
        ("③ PLG（Product-Led Growth）の浸透",
         "フリーミアム→有料転換モデルが主流に。SMB獲得コストを大幅削減し、国内スタートアップも採用拡大。",
         "Medium"),
        ("④ セキュリティ・コンプライアンス強化",
         "政府のクラウド移行方針とともにISMAP認証・ゼロトラスト対応が選定条件に浮上。",
         "Medium"),
    ]

    priority_color = {"High": ORANGE, "Medium": BLUE}
    tw = Inches(5.9)
    th = Inches(1.85)
    gap = Inches(0.18)
    positions = [
        (Inches(0.3),  Inches(1.35)),
        (Inches(6.55), Inches(1.35)),
        (Inches(0.3),  Inches(3.38)),
        (Inches(6.55), Inches(3.38)),
    ]

    for (tx, ty), (title, body, prio) in zip(positions, trends):
        add_rect(sl, tx, ty, tw, th, WHITE, border=True, border_color=LGRAY)
        # 左色帯
        pcol = priority_color.get(prio, BLUE)
        add_rect(sl, tx, ty, Inches(0.18), th, pcol)
        # タイトル
        add_text(sl, title, tx + Inches(0.25), ty + Inches(0.1),
                 tw - Inches(0.3), Inches(0.42),
                 font_size=11, bold=True, color=NAVY)
        # 優先度バッジ
        add_rect(sl, tx + tw - Inches(0.9), ty + Inches(0.08),
                 Inches(0.82), Inches(0.28), pcol)
        add_text(sl, prio, tx + tw - Inches(0.88), ty + Inches(0.09),
                 Inches(0.8), Inches(0.25),
                 font_size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 本文
        add_text(sl, body, tx + Inches(0.25), ty + Inches(0.58),
                 tw - Inches(0.3), Inches(1.15),
                 font_size=9.5, color=BLACK)

    footer(sl)

# ============================================================
# スライド6: ビジネス示唆
# ============================================================
def slide_implications(prs):
    sl = blank_slide(prs)
    header_bar(sl, "ビジネス示唆 — 参入・投資の優先アクション",
               "Business Implications & Recommended Actions")

    implications = [
        ("High", "バーティカルSaaSへの参入検討",
         "汎用SaaSが成熟しつつある現在、業界特化型は競争が少なく高ARRを狙える。医療・建設・製造を優先領域として検討。"),
        ("High", "AI機能の早期組み込み",
         "LLM統合は今や顧客の期待値。競合に対する差別化として、既存プロダクトへのCopilot的機能追加を最優先で投資。"),
        ("Medium", "PLGモデルへの転換検討",
         "SMB市場獲得にはフリーミアム→有料転換モデルが有効。営業コスト削減とロングテール顧客獲得を同時に実現。"),
        ("Medium", "ISMAP・セキュリティ認証取得",
         "政府・自治体案件を狙う場合、ISMAP登録は必須要件。認証取得で公共市場への参入障壁を突破できる。"),
    ]

    priority_color = {"High": ORANGE, "Medium": BLUE}
    row_h = Inches(1.38)
    row_gap = Inches(0.14)
    sx = Inches(0.35)
    sy = Inches(1.38)

    for i, (prio, title, body) in enumerate(implications):
        ry = sy + i * (row_h + row_gap)
        pcol = priority_color[prio]
        # 行背景
        bg = LGRAY if i % 2 == 0 else WHITE
        add_rect(sl, sx, ry, Inches(12.6), row_h, bg, border=True, border_color=LGRAY)
        # 優先度帯
        add_rect(sl, sx, ry, Inches(0.85), row_h, pcol)
        add_text(sl, prio, sx + Inches(0.02), ry + Inches(0.55),
                 Inches(0.8), Inches(0.38),
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # タイトル
        add_text(sl, f"▶ {title}", sx + Inches(0.95), ry + Inches(0.1),
                 Inches(3.8), Inches(0.45),
                 font_size=12, bold=True, color=NAVY)
        # 本文
        add_text(sl, body, sx + Inches(0.95), ry + Inches(0.58),
                 Inches(11.3), Inches(0.7),
                 font_size=9.5, color=BLACK)

    footer(sl)

# ============================================================
# メイン
# ============================================================
OUT = r"C:\Users\tshibasaki\Desktop\etc\work\task\05_CONTENT\drafts\saas_market_2025\SaaS市場調査レポート_2025_v1.pptx"

prs = new_prs()
slide_cover(prs)
slide_exec_summary(prs)
slide_market_size(prs)
slide_players(prs)
slide_trends(prs)
slide_implications(prs)
prs.save(OUT)
print(f"Done: {OUT}")
