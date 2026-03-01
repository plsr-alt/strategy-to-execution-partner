"""
SaaS市場調査スライド v2 — BCGスタイル完全再現
・白背景
・タイトル＝結論メッセージ（文章形式）
・グリーン系アクセント
・出所: 左下
・グッドプラクティス風バッジ
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ===== BCGカラーパレット =====
GREEN_DARK  = RGBColor(0x1A, 0x5C, 0x2A)   # 濃緑 (タイトル・強調)
GREEN_MID   = RGBColor(0x2E, 0x8B, 0x45)   # 中緑 (バー・ヘッダー)
GREEN_LIGHT = RGBColor(0xA8, 0xD5, 0xB5)   # 薄緑 (サブバー・背景帯)
GREEN_PALE  = RGBColor(0xE8, 0xF5, 0xEC)   # 極薄緑 (背景)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_DARK   = RGBColor(0x44, 0x44, 0x44)
GRAY_MID    = RGBColor(0x88, 0x88, 0x88)
GRAY_LIGHT  = RGBColor(0xE0, 0xE0, 0xE0)
RED_ACCENT  = RGBColor(0xCC, 0x33, 0x33)   # 強調サークル用

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill, border_color=None, border_pt=0.75):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=11, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, name="Yu Gothic"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = name
    return tb

def add_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    """水平・垂直ライン"""
    from pptx.util import Emu
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn

# ============================================================
# 共通パーツ
# ============================================================
def bcg_title(slide, message, sub="", badge=""):
    """
    BCGスタイルタイトル:
    - バッジ（グッドプラクティス等）
    - 結論メッセージ（太字・濃緑）
    - サブタイトル（グレー）
    - 下線
    """
    y = Inches(0.25)
    if badge:
        rect(slide, Inches(0.4), y, Inches(1.5), Inches(0.26), GREEN_MID)
        txt(slide, badge,
            Inches(0.42), y + Pt(2), Inches(1.46), Inches(0.24),
            size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.3)

    txt(slide, message,
        Inches(0.4), y, Inches(12.5), Inches(0.9),
        size=18, bold=True, color=GREEN_DARK)
    y += Inches(0.92)

    if sub:
        txt(slide, sub,
            Inches(0.4), y, Inches(12.0), Inches(0.3),
            size=10, color=GRAY_DARK)
        y += Inches(0.32)

    # 区切り線
    add_line(slide, Inches(0.4), y, Inches(12.9), y, GREEN_MID, 1.5)
    return y + Inches(0.08)

def source_note(slide, text="出所：Gartner調査 / 各社公開情報 / BCG分析"):
    y = SLIDE_H - Inches(0.28)
    add_line(slide, Inches(0.4), y, Inches(12.9), y, GRAY_LIGHT, 0.5)
    txt(slide, text,
        Inches(0.4), y + Pt(2), Inches(10), Inches(0.24),
        size=7, color=GRAY_MID)

def page_num(slide, n):
    txt(slide, str(n),
        SLIDE_W - Inches(0.5), SLIDE_H - Inches(0.28),
        Inches(0.4), Inches(0.24),
        size=8, color=GRAY_MID, align=PP_ALIGN.RIGHT)

# ============================================================
# Slide 1: 表紙
# ============================================================
def slide_cover(prs):
    sl = blank_slide(prs)
    # 白背景＋左グリーン帯
    rect(sl, 0, 0, Inches(0.6), SLIDE_H, GREEN_DARK)
    # タイトル
    txt(sl, "日本のSaaS市場\n2025年 市場調査レポート",
        Inches(1.0), Inches(1.8), Inches(10.0), Inches(2.5),
        size=32, bold=True, color=GREEN_DARK, align=PP_ALIGN.LEFT)
    # 区切り線
    add_line(sl, Inches(1.0), Inches(4.35), Inches(8.0), Inches(4.35), GREEN_MID, 2)
    # サブ
    txt(sl, "市場規模 / 主要プレイヤー / トレンド / ビジネス示唆",
        Inches(1.0), Inches(4.5), Inches(10.0), Inches(0.5),
        size=13, color=GRAY_DARK)
    txt(sl, "2026年2月",
        Inches(1.0), Inches(5.2), Inches(4.0), Inches(0.35),
        size=10, color=GRAY_MID)
    source_note(sl, "")

# ============================================================
# Slide 2: エグゼクティブサマリー
# ============================================================
def slide_exec_summary(prs):
    sl = blank_slide(prs)
    content_y = bcg_title(
        sl,
        "日本のSaaS市場はDX推進を背景に急拡大しており、AI統合・業界特化が次の成長波となる",
        sub="エグゼクティブサマリー｜Key Findings"
    )

    # 3カラムボックス
    findings = [
        ("01  市場規模",
         "世界SaaS市場\n1,030億ドル（2020）\n→ 1,450億ドル（2022）\nCAGR ≈ 11%\n\n日本市場は2024年に\n1兆円超と推計"),
        ("02  主要プレイヤー",
         "Tier1（グローバル）\n・Salesforce — CRM\n・Microsoft — 365/Teams\n\nTier2\n・Google Workspace\n・Adobe Experience\n\n国内勢\n・マネーフォワード\n・freee / Sansan"),
        ("03  ビジネス示唆",
         "【High】\n・バーティカルSaaSへの参入\n・AI機能の早期統合\n\n【Medium】\n・PLGモデル採用\n・ISMAP認証取得"),
    ]

    bw = Inches(3.9)
    bh = SLIDE_H - content_y - Inches(0.5)
    gap = Inches(0.25)
    bx = Inches(0.35)

    for i, (hd, body) in enumerate(findings):
        x = bx + i * (bw + gap)
        # ヘッダー帯
        rect(sl, x, content_y, bw, Inches(0.38), GREEN_MID)
        txt(sl, hd, x + Inches(0.1), content_y + Pt(4),
            bw - Inches(0.2), Inches(0.32),
            size=10, bold=True, color=WHITE)
        # 本文エリア
        rect(sl, x, content_y + Inches(0.38), bw, bh - Inches(0.38),
             WHITE, border_color=GRAY_LIGHT)
        txt(sl, body,
            x + Inches(0.15), content_y + Inches(0.5),
            bw - Inches(0.3), bh - Inches(0.6),
            size=9.5, color=BLACK)

    source_note(sl)
    page_num(sl, 2)

# ============================================================
# Slide 3: 市場規模（BCG棒グラフ風）
# ============================================================
def slide_market_size(prs):
    sl = blank_slide(prs)
    content_y = bcg_title(
        sl,
        "世界SaaS市場は2020年比約40%増の1,450億ドルへ急拡大、CAGR11%の成長が継続している",
        sub="市場規模推移｜Global SaaS Market Size"
    )

    # --- 棒グラフ（図形で表現） ---
    chart_x  = Inches(0.8)
    chart_bx = Inches(1.6)    # バー開始X
    chart_y  = content_y + Inches(0.2)
    max_h    = Inches(3.6)
    max_val  = 2100
    bar_w    = Inches(1.3)
    gap      = Inches(0.9)

    bars = [
        ("2020年", 1030, GREEN_LIGHT, "1,030億ドル"),
        ("2022年", 1450, GREEN_MID,   "1,450億ドル"),
        ("2025年\n(予測)", 2100, GREEN_DARK, "~2,100億ドル"),
    ]

    # Y軸ライン
    add_line(sl, chart_x, chart_y, chart_x, chart_y + max_h + Inches(0.05), GRAY_DARK, 1)
    # X軸ライン
    add_line(sl, chart_x, chart_y + max_h, chart_x + Inches(6.5), chart_y + max_h, GRAY_DARK, 1)

    for i, (yr, val, col, label) in enumerate(bars):
        bh_val = max_h * val / max_val
        bx = chart_bx + i * (bar_w + gap)
        by = chart_y + max_h - bh_val
        # バー
        rect(sl, bx, by, bar_w, bh_val, col)
        # 数値ラベル（バー上）
        txt(sl, label, bx - Inches(0.1), by - Inches(0.42),
            bar_w + Inches(0.2), Inches(0.38),
            size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
        # X軸ラベル
        txt(sl, yr, bx, chart_y + max_h + Inches(0.07),
            bar_w, Inches(0.45),
            size=9, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # CAGR矢印注釈
    txt(sl, "CAGR\n≈ 11%",
        Inches(4.5), chart_y + Inches(0.4),
        Inches(1.2), Inches(0.65),
        size=11, bold=True, color=RED_ACCENT, align=PP_ALIGN.CENTER)

    # --- 右側: 日本市場ボックス ---
    rx = Inches(8.2)
    ry = content_y + Inches(0.2)
    rw = Inches(4.8)
    rh = Inches(4.5)
    rect(sl, rx, ry, rw, rh, GREEN_PALE, border_color=GREEN_MID)
    rect(sl, rx, ry, rw, Inches(0.38), GREEN_MID)
    txt(sl, "日本市場 フォーカス",
        rx + Inches(0.1), ry + Pt(4), rw - Inches(0.2), Inches(0.3),
        size=10, bold=True, color=WHITE)

    jp_items = [
        ("規模推定",    "2024年に国内SaaS市場が\n1兆円超に到達と推計"),
        ("成長背景",    "政府DX推進・クラウド\n移行方針が追い風"),
        ("普及率",      "国内企業の64%以上が\nSaaS含むクラウド導入済"),
        ("今後の鍵",    "中小企業への浸透が\n次の成長ドライバー"),
    ]
    item_y = ry + Inches(0.5)
    for label, desc in jp_items:
        rect(sl, rx + Inches(0.15), item_y, Inches(0.06), Inches(0.55), GREEN_MID)
        txt(sl, label, rx + Inches(0.3), item_y,
            Inches(1.2), Inches(0.3), size=8, bold=True, color=GREEN_DARK)
        txt(sl, desc,  rx + Inches(0.3), item_y + Inches(0.28),
            rw - Inches(0.5), Inches(0.35), size=8.5, color=BLACK)
        item_y += Inches(0.95)

    source_note(sl, "出所：Gartner調査（2022）/ スマートキャンプ SaaS業界レポート / 各種公開情報")
    page_num(sl, 3)

# ============================================================
# Slide 4: 主要プレイヤー
# ============================================================
def slide_players(prs):
    sl = blank_slide(prs)
    content_y = bcg_title(
        sl,
        "SalesforceとMicrosoftがグローバルをリード、国内ではマネーフォワード・freeeが急成長中",
        sub="主要プレイヤー分類｜Key Players by Tier"
    )

    tiers = [
        ("Tier 1  グローバルリーダー", GREEN_DARK, [
            ("Salesforce",
             "CRM・営業支援のデファクトスタンダード。国内大企業の採用率が高く、Einstein AIで差別化を強化。"),
            ("Microsoft",
             "Microsoft 365・Teams・Power Platformで業務全体をカバー。Copilot AI統合で圧倒的優位。"),
        ]),
        ("Tier 2  メジャープレイヤー", GREEN_MID, [
            ("Google Workspace",
             "旧G Suite。SMB市場を中心に拡大。低コストと使いやすさが強み。Gemini AI統合が進行中。"),
            ("Adobe",
             "Experience Cloudでマーケティング・DX領域に特化。クリエイティブ×データ分析が武器。"),
        ]),
        ("Tier 3  国内注目プレイヤー", GREEN_LIGHT, [
            ("マネーフォワード",
             "クラウド会計・給与で中小企業向けに急成長。バックオフィス全般のSaaS化を推進。"),
            ("freee / Sansan",
             "freeeは会計・HR、SansanはCRM×名刺管理で差別化。国内SMBのDX浸透を牽引。"),
        ]),
    ]

    col_w = Inches(4.1)
    col_h = SLIDE_H - content_y - Inches(0.45)
    gap   = Inches(0.25)
    sx    = Inches(0.35)

    for ci, (tier_name, col, items) in enumerate(tiers):
        cx = sx + ci * (col_w + gap)
        # ヘッダー
        rect(sl, cx, content_y, col_w, Inches(0.42), col)
        header_text_color = BLACK if col == GREEN_LIGHT else WHITE
        txt(sl, tier_name,
            cx + Inches(0.12), content_y + Pt(5),
            col_w - Inches(0.2), Inches(0.35),
            size=9.5, bold=True, color=header_text_color)

        # カード2枚
        card_h = (col_h - Inches(0.42) - Inches(0.12)) / 2
        for ri, (name, desc) in enumerate(items):
            ry = content_y + Inches(0.42) + ri * (card_h + Inches(0.06))
            rect(sl, cx, ry, col_w, card_h, WHITE, border_color=GRAY_LIGHT)
            # 左アクセントライン
            rect(sl, cx, ry, Inches(0.06), card_h, col)
            # 企業名
            txt(sl, name,
                cx + Inches(0.15), ry + Inches(0.1),
                col_w - Inches(0.2), Inches(0.38),
                size=12, bold=True, color=GREEN_DARK)
            # 区切り
            add_line(sl, cx + Inches(0.15), ry + Inches(0.52),
                     cx + col_w - Inches(0.15), ry + Inches(0.52),
                     GRAY_LIGHT, 0.5)
            # 説明
            txt(sl, desc,
                cx + Inches(0.15), ry + Inches(0.6),
                col_w - Inches(0.25), card_h - Inches(0.7),
                size=9, color=GRAY_DARK)

    source_note(sl, "出所：各社公開情報 / Salesforce Ventures調査")
    page_num(sl, 4)

# ============================================================
# Slide 5: トレンド
# ============================================================
def slide_trends(prs):
    sl = blank_slide(prs)
    content_y = bcg_title(
        sl,
        "AI統合・バーティカル特化・PLGの3潮流が2025年のSaaS市場の次なる成長エンジンとなる",
        sub="主要トレンド｜Key Trends 2025",
        badge="調査サマリー"
    )

    trends = [
        ("① AI統合SaaSの台頭",
         "ChatGPT等のLLMがSaaS製品に組み込まれ業務自動化が急加速。"
         "Salesforce Einstein・Microsoft Copilotが先行し、"
         "「AI前提」の製品設計が業界標準になりつつある。",
         "High", True),
        ("② バーティカルSaaS（業界特化）",
         "医療・建設・製造など業界特化型SaaSが急成長。"
         "汎用SaaSより解約率が低く高ARRを実現。"
         "国内でも医療DX・建設DX領域で新興プレイヤーが台頭。",
         "High", True),
        ("③ PLG（Product-Led Growth）の浸透",
         "フリーミアム→有料転換モデルが主流に。"
         "エンドユーザーが試して採用するボトムアップ型で、"
         "SMB獲得コストを大幅削減。国内SaaSも積極採用中。",
         "Medium", False),
        ("④ セキュリティ・コンプライアンス強化",
         "政府のクラウド移行方針とともにISMAP認証・"
         "ゼロトラスト対応が選定条件に浮上。"
         "公共・金融分野への参入には必須要件となってきた。",
         "Medium", False),
    ]

    tw = Inches(6.0)
    th = Inches(2.35)
    gap = Inches(0.2)
    positions = [
        (Inches(0.35), content_y),
        (Inches(6.95), content_y),
        (Inches(0.35), content_y + th + gap),
        (Inches(6.95), content_y + th + gap),
    ]
    prio_color = {"High": RED_ACCENT, "Medium": GREEN_MID}

    for (tx, ty), (title, body, prio, is_high) in zip(positions, trends):
        # 外枠
        rect(sl, tx, ty, tw, th, WHITE, border_color=GRAY_LIGHT)
        # 上帯（Highのみ濃緑、Mediumは薄緑）
        header_col = GREEN_DARK if is_high else GREEN_LIGHT
        header_txt = WHITE if is_high else BLACK
        rect(sl, tx, ty, tw, Inches(0.38), header_col)
        # 優先度バッジ
        pc = prio_color[prio]
        rect(sl, tx + tw - Inches(0.85), ty + Inches(0.05),
             Inches(0.78), Inches(0.26), pc)
        txt(sl, prio,
            tx + tw - Inches(0.83), ty + Inches(0.06),
            Inches(0.74), Inches(0.24),
            size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # タイトル
        txt(sl, title,
            tx + Inches(0.12), ty + Inches(0.05),
            tw - Inches(1.0), Inches(0.3),
            size=10.5, bold=True, color=header_txt)
        # 本文
        txt(sl, body,
            tx + Inches(0.12), ty + Inches(0.44),
            tw - Inches(0.22), th - Inches(0.54),
            size=9.5, color=BLACK)

    source_note(sl, "出所：各種業界調査 / FastGrow / Salesforce Ventures / 公開情報")
    page_num(sl, 5)

# ============================================================
# Slide 6: ビジネス示唆
# ============================================================
def slide_implications(prs):
    sl = blank_slide(prs)
    content_y = bcg_title(
        sl,
        "参入優先度は「バーティカルSaaS×AI機能統合」が最上位、まずPLGモデルの設計から着手すべき",
        sub="ビジネス示唆｜Recommended Actions",
        badge="示唆・提言"
    )

    rows = [
        ("High",   "バーティカルSaaSへの参入検討",
         "汎用SaaSが成熟しつつある現在、業界特化型は競争が少なく高ARRを狙える。医療・建設・製造を優先領域として検討。参入コストを抑えるため既存SaaS APIとの連携設計が鍵。"),
        ("High",   "AI機能の早期組み込み",
         "LLM統合は今や顧客の期待値。競合に対する差別化として既存プロダクトへのCopilot的機能追加を最優先で投資。PoC→3ヶ月での本番展開を目安に計画を立てる。"),
        ("Medium", "PLGモデルへの転換検討",
         "SMB市場獲得にはフリーミアム→有料転換モデルが有効。営業コスト削減とロングテール顧客獲得を同時に実現。Onboarding体験設計がPLG成否の核心。"),
        ("Medium", "ISMAP・セキュリティ認証取得",
         "政府・自治体案件を狙う場合、ISMAP登録は必須要件。認証取得で公共市場への参入障壁を突破できる。申請〜取得まで12〜18ヶ月を想定した先行投資が必要。"),
    ]

    rh   = Inches(1.25)
    rgap = Inches(0.1)
    sx   = Inches(0.35)
    rw   = Inches(12.6)
    prio_col = {"High": RED_ACCENT, "Medium": GREEN_MID}

    for i, (prio, title, body) in enumerate(rows):
        ry = content_y + i * (rh + rgap)
        bg = GREEN_PALE if i % 2 == 0 else WHITE
        rect(sl, sx, ry, rw, rh, bg, border_color=GRAY_LIGHT, border_pt=0.5)

        # 優先度ラベル列
        pc = prio_col[prio]
        rect(sl, sx, ry, Inches(0.9), rh, pc)
        txt(sl, prio,
            sx + Inches(0.02), ry + Inches(0.45),
            Inches(0.86), Inches(0.35),
            size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # タイトル
        txt(sl, f"▶  {title}",
            sx + Inches(1.0), ry + Inches(0.1),
            Inches(4.0), Inches(0.38),
            size=11.5, bold=True, color=GREEN_DARK)

        # 本文
        txt(sl, body,
            sx + Inches(1.0), ry + Inches(0.52),
            rw - Inches(1.15), Inches(0.65),
            size=9, color=GRAY_DARK)

    source_note(sl, "出所：各種市場調査 / BCG分析フレームワーク参照")
    page_num(sl, 6)

# ============================================================
# 実行
# ============================================================
OUT = r"C:\Users\tshibasaki\Desktop\etc\work\task\05_CONTENT\drafts\saas_market_2025\SaaS市場調査レポート_2025_v2.pptx"

prs = new_prs()
slide_cover(prs)
slide_exec_summary(prs)
slide_market_size(prs)
slide_players(prs)
slide_trends(prs)
slide_implications(prs)
prs.save(OUT)
print(f"Done: {OUT}")
