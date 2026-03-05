#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
コメダ様 FSx選択肢メリデメ PPT生成スクリプト v2
コンサルクオリティ（McKinsey/BCG/Accenture風）

デザイン原則:
- 16:9 ワイドスクリーン
- 2-3色 + グレー階調
- Assertion-Evidence 構造（タイトルは主張文）
- ホワイトスペース30-40%
- フォント3段階（タイトル/本文/注釈）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from datetime import datetime
import os
import platform

# ============================================================
# 定数定義
# ============================================================

# --- カラーパレット ---
COLOR_NAVY     = RGBColor(0x1F, 0x49, 0x7D)  # プライマリ（濃紺）
COLOR_BLUE     = RGBColor(0x2B, 0x57, 0x9A)  # セカンダリ（青）
COLOR_ACCENT   = RGBColor(0x00, 0x78, 0xD4)  # アクセント（ティール）
COLOR_RED      = RGBColor(0xC0, 0x00, 0x00)  # 警告
COLOR_GREEN    = RGBColor(0x2E, 0x7D, 0x32)  # 成功
COLOR_TEXT     = RGBColor(0x33, 0x33, 0x33)  # 本文テキスト
COLOR_SUBTEXT  = RGBColor(0x66, 0x66, 0x66)  # 注釈・出典
COLOR_BG_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)  # 背景グレー
COLOR_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BORDER   = RGBColor(0xD9, 0xD9, 0xD9)  # 罫線
COLOR_LIGHT_NAVY = RGBColor(0xE8, 0xEE, 0xF4)  # 薄い紺（セクション背景）

# --- フォントサイズ ---
FONT_TITLE     = Pt(20)
FONT_SECTION   = Pt(28)
FONT_SUBHEAD   = Pt(14)
FONT_BODY      = Pt(11)
FONT_BODY_L    = Pt(12)
FONT_NOTE      = Pt(8)
FONT_KPI       = Pt(32)
FONT_PAGE_NUM  = Pt(8)

# --- フォント ---
FONT_JP = "メイリオ"
FONT_EN = "Segoe UI"

# --- レイアウト ---
MARGIN_LEFT    = Inches(0.75)
MARGIN_RIGHT   = Inches(0.75)
MARGIN_TOP     = Inches(0.5)
SLIDE_WIDTH    = Inches(13.333)
SLIDE_HEIGHT   = Inches(7.5)

TITLE_LEFT     = MARGIN_LEFT
TITLE_TOP      = Inches(0.4)
TITLE_WIDTH    = Inches(11.8)
TITLE_HEIGHT   = Inches(0.65)

CONTENT_LEFT   = MARGIN_LEFT
CONTENT_TOP    = Inches(1.4)
CONTENT_WIDTH  = Inches(11.8)
CONTENT_HEIGHT = Inches(5.5)


# ============================================================
# ヘルパー関数
# ============================================================

def set_font(run, size=FONT_BODY, bold=False, color=COLOR_TEXT, name_jp=FONT_JP, name_en=FONT_EN):
    """フォントを一括設定（日本語フォント対応）"""
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name_en
    # 東アジアフォント設定
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name_jp)


def add_textbox(slide, left, top, width, height, text, size=FONT_BODY,
                bold=False, color=COLOR_TEXT, alignment=PP_ALIGN.LEFT,
                word_wrap=True):
    """テキストボックスを追加して Run を返す"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return txBox, tf


def add_title_bar(slide, text):
    """コンサルスタイルのタイトルバー（主張文 + アクセントライン）"""
    txBox, tf = add_textbox(
        slide, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        text, size=FONT_TITLE, bold=True, color=COLOR_NAVY
    )
    # アクセントライン
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        TITLE_LEFT, TITLE_TOP + TITLE_HEIGHT + Inches(0.05),
        TITLE_WIDTH, Inches(0.025)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()
    return txBox


def add_page_number(slide, number, total=None):
    """ページ番号を右下に追加"""
    if total:
        text = f"{number} / {total}"
    else:
        text = str(number)
    add_textbox(
        slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
        text, size=FONT_PAGE_NUM, color=COLOR_SUBTEXT,
        alignment=PP_ALIGN.RIGHT
    )


def add_source_note(slide, text, top=Inches(6.8)):
    """出典・注釈を左下に追加"""
    add_textbox(
        slide, MARGIN_LEFT, top, Inches(10.0), Inches(0.3),
        text, size=FONT_NOTE, color=COLOR_SUBTEXT
    )


def add_footer_bar(slide):
    """フッターバー（薄いライン）"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.2),
        SLIDE_WIDTH, Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BORDER
    line.line.fill.background()


def _set_cell_color(cell, color):
    """セル背景色を設定"""
    tcPr = cell._tc.get_or_add_tcPr()
    # 既存のsolidFillを削除
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))


def _set_cell_border(cell, border_pos, color, width_emu=12700):
    """セルの罫線を設定"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    border_map = {
        'top': 'a:lnT', 'bottom': 'a:lnB',
        'left': 'a:lnL', 'right': 'a:lnR'
    }
    tag = qn(border_map[border_pos])
    # 既存を削除
    for old in tcPr.findall(tag):
        tcPr.remove(old)
    ln = etree.SubElement(tcPr, tag)
    ln.set('w', str(width_emu))
    solidFill = etree.SubElement(ln, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))


def set_cell_margins(cell, top=45720, bottom=45720, left=91440, right=91440):
    """セル内余白を設定"""
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.set('marT', str(top))
    tcPr.set('marB', str(bottom))
    tcPr.set('marL', str(left))
    tcPr.set('marR', str(right))


def create_table(slide, data, left=CONTENT_LEFT, top=CONTENT_TOP,
                 width=CONTENT_WIDTH, height=None, col_widths=None,
                 header_color=COLOR_NAVY):
    """コンサルクオリティのテーブルを作成"""
    rows = len(data)
    cols = len(data[0])
    if height is None:
        height = Inches(min(rows * 0.55, 5.5))

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 列幅
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        col_w = int(width / cols)
        for i in range(cols):
            table.columns[i].width = col_w

    # ヘッダー行
    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = data[0][j]
        set_font(run, size=FONT_BODY, bold=True, color=COLOR_WHITE)
        _set_cell_color(cell, header_color)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_cell_margins(cell)

    # データ行
    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = data[i][j]
            set_font(run, size=Pt(10), color=COLOR_TEXT)
            # 交互色
            if i % 2 == 0:
                _set_cell_color(cell, COLOR_BG_GRAY)
            else:
                _set_cell_color(cell, COLOR_WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            set_cell_margins(cell)

    # 罫線
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            for pos in ['top', 'bottom', 'left', 'right']:
                _set_cell_border(cell, pos, COLOR_BORDER)

    return table


def add_bullet_content(slide, sections, left=CONTENT_LEFT, top=CONTENT_TOP,
                       width=CONTENT_WIDTH, line_spacing=1.3):
    """
    構造化された箇条書きコンテンツを追加
    sections: [{"heading": str, "items": [str], "color": RGBColor (optional)}]
    """
    current_top = top
    for section in sections:
        heading = section.get("heading", "")
        items = section.get("items", [])
        h_color = section.get("color", COLOR_NAVY)

        if heading:
            _, tf = add_textbox(
                slide, left, current_top, width, Inches(0.35),
                heading, size=FONT_SUBHEAD, bold=True, color=h_color
            )
            current_top += Inches(0.35)

        for item in items:
            _, tf = add_textbox(
                slide, left + Inches(0.3), current_top, width - Inches(0.3), Inches(0.3),
                f"• {item}", size=FONT_BODY, color=COLOR_TEXT
            )
            current_top += Inches(0.28)

        current_top += Inches(0.15)  # セクション間スペース

    return current_top


def add_two_column_content(slide, left_sections, right_sections):
    """2カラムレイアウトでコンテンツを配置"""
    col_width = Inches(5.5)
    left_x = CONTENT_LEFT
    right_x = CONTENT_LEFT + col_width + Inches(0.5)

    add_bullet_content(slide, left_sections, left=left_x, width=col_width)
    add_bullet_content(slide, right_sections, left=right_x, width=col_width)


# ============================================================
# スライド生成
# ============================================================

def slide_01_title(prs):
    """スライド1: 表紙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 上半分: 濃紺背景
    top_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(4.5)
    )
    top_bg.fill.solid()
    top_bg.fill.fore_color.rgb = COLOR_NAVY
    top_bg.line.fill.background()

    # アクセントライン
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(0.06)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLOR_ACCENT
    accent_line.line.fill.background()

    # タイトル
    add_textbox(
        slide, Inches(1.2), Inches(1.5), Inches(11.0), Inches(1.2),
        "AWS ファイルストレージ選択肢\nメリット・デメリット比較",
        size=Pt(36), bold=True, color=COLOR_WHITE,
        alignment=PP_ALIGN.LEFT
    )

    # サブタイトル
    add_textbox(
        slide, Inches(1.2), Inches(3.2), Inches(11.0), Inches(0.8),
        "FSx for Windows File Server / FSx for NetApp ONTAP / Qumulo",
        size=Pt(16), color=RGBColor(0xB0, 0xC4, 0xDE),
        alignment=PP_ALIGN.LEFT
    )

    # 下部情報
    add_textbox(
        slide, Inches(1.2), Inches(5.2), Inches(5.0), Inches(0.5),
        "コメダ様向け | ファイルサーバークラウド移行支援",
        size=FONT_BODY_L, color=COLOR_SUBTEXT
    )
    add_textbox(
        slide, Inches(1.2), Inches(5.7), Inches(5.0), Inches(0.5),
        f"作成日: {datetime.now().strftime('%Y年%m月%d日')}",
        size=FONT_BODY, color=COLOR_SUBTEXT
    )
    add_textbox(
        slide, Inches(1.2), Inches(6.1), Inches(5.0), Inches(0.5),
        "CONFIDENTIAL",
        size=FONT_NOTE, bold=True, color=COLOR_SUBTEXT
    )


def slide_02_executive_summary(prs):
    """スライド2: エグゼクティブサマリ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "3つの選択肢はいずれもAWS上で利用可能だが、特性が大きく異なる")
    add_page_number(slide, 2, 9)
    add_footer_bar(slide)

    sections = [
        {
            "heading": "背景",
            "color": COLOR_NAVY,
            "items": [
                "現行オンプレミス Windows ファイルサーバー（約50TB）のAWSクラウド移行を検討中",
                "移行後の運用体制・コスト・学習コストのバランスが重要な判断基準",
            ]
        },
        {
            "heading": "選択肢",
            "color": COLOR_NAVY,
            "items": [
                "FSx for Windows File Server — Windows運用をそのまま継承。SMB専用。完全マネージド",
                "FSx for NetApp ONTAP — マルチプロトコル（NFS+SMB+iSCSI）。階層化ストレージで柔軟。完全マネージド",
                "Qumulo on AWS — EC2ベースのソフトウェアデプロイ型。PB級スケーラビリティ。セルフマネージド",
            ]
        },
        {
            "heading": "推奨",
            "color": COLOR_GREEN,
            "items": [
                "50TBスケールかつWindows運用継承を重視する場合、FSx for Windows File Server（HDD構成）が最有力",
                "NFS要件やデータ圧縮効果が見込める場合、FSx for NetApp ONTAP も有力な候補",
            ]
        },
    ]
    add_bullet_content(slide, sections)


def slide_03_comparison_matrix(prs):
    """スライド3: 総合比較マトリクス"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "FSx for Windows が運用継承性とコストのバランスで最も優位")
    add_page_number(slide, 3, 9)
    add_footer_bar(slide)

    data = [
        ["比較項目", "FSx for Windows\nFile Server", "FSx for NetApp\nONTAP", "Qumulo\non AWS"],
        ["サービス種別", "完全マネージド\n（AWS純正）", "完全マネージド\n（AWS純正）", "セルフマネージド\n（EC2ベース）"],
        ["対応プロトコル", "SMB", "NFS + SMB + iSCSI", "NFS + SMB + S3"],
        ["概算費用（50TB/年）\n※東京リージョン", "約250〜270万円\n（HDD / Single-AZ）", "約250〜420万円\n（SSD+Capacity Pool）", "要個別見積\n（EC2+ライセンス費）"],
        ["運用難易度\n（Windowsスキル前提）", "低〜中 ◎\nWindows管理者に最も馴染む", "中 ○\nONTAP管理知識が必要", "中〜高 △\nEC2インフラ管理が必要"],
        ["バックアップ", "AWS Backup統合\n自動日次バックアップ", "ONTAP Snapshot\n+ SnapMirror", "Qumulo Snapshot\n+ 独自設計が必要"],
        ["AD連携", "◎ 必須・完全統合", "○ SMB利用時対応", "○ SMB利用時対応"],
        ["データ重複排除\n/ 圧縮", "○ Windows機能", "◎ ONTAP標準\n（高い圧縮効果）", "○ 対応"],
        ["スケーラビリティ", "最大 64TB(SSD)\n/ 512TB(HDD)", "SSD: 192TB\nCapacity Pool: 無制限", "PB級まで拡張可能"],
    ]

    col_widths = [Inches(2.2), Inches(3.2), Inches(3.2), Inches(3.2)]
    create_table(slide, data, top=Inches(1.3), height=Inches(5.5), col_widths=col_widths)

    add_source_note(slide, "※ 費用は2025年時点のAWS公開価格に基づく概算。リザーブド価格・圧縮効果により変動。最終見積はAWS料金計算ツールで確認推奨")


def slide_04_fsx_windows(prs):
    """スライド4: FSx for Windows File Server 詳細"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "選択肢1: FSx for Windows — 既存運用の継承性が最大の強み")
    add_page_number(slide, 4, 9)
    add_footer_bar(slide)

    left_sections = [
        {
            "heading": "✅ メリット",
            "color": COLOR_GREEN,
            "items": [
                "既存Windowsファイルサーバーの運用方法をそのまま継承可能",
                "MMC / PowerShell など馴染みのある管理ツールで操作",
                "Active Directory との完全統合（細かいACL制御）",
                "AWS Backup による自動日次バックアップ",
                "DFS名前空間・データ重複排除をサポート",
                "AWS完全マネージド — インフラ管理不要",
            ]
        },
    ]

    right_sections = [
        {
            "heading": "⚠ デメリット / 注意点",
            "color": COLOR_RED,
            "items": [
                "SMBプロトコルのみ（NFS非対応）",
                "AWS側の監視設計・IAM管理は別途構築が必要",
                "リストアは新規ファイルシステムとして復元（上書きリストア不可）",
                "SSD構成は高額 → HDD構成を推奨（50TBスケール）",
                "3選択肢中、費用面ではFSx ONTAPと同等水準",
            ]
        },
    ]

    add_two_column_content(slide, left_sections, right_sections)

    # コスト概算カード
    cost_top = Inches(5.0)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        CONTENT_LEFT, cost_top, Inches(11.8), Inches(1.3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_NAVY
    card.line.color.rgb = COLOR_BORDER
    card.line.width = Pt(0.5)

    add_textbox(
        slide, CONTENT_LEFT + Inches(0.3), cost_top + Inches(0.15),
        Inches(11.2), Inches(0.3),
        "コスト概算（50TB / 東京リージョン / Single-AZ / HDD構成）",
        size=FONT_BODY, bold=True, color=COLOR_NAVY
    )
    add_textbox(
        slide, CONTENT_LEFT + Inches(0.3), cost_top + Inches(0.5),
        Inches(11.2), Inches(0.6),
        "HDD: $0.025/GB/月 × 50TB = $1,250/月 + スループット(32MB/s): 約$176/月 → 年間約250〜270万円（1$=150円換算）",
        size=Pt(10), color=COLOR_TEXT
    )

    add_source_note(slide, "出典: AWS公式料金ページ（ap-northeast-1） | ※ Multi-AZ構成の場合、費用は約2倍")


def slide_05_fsx_ontap(prs):
    """スライド5: FSx for NetApp ONTAP 詳細"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "選択肢2: FSx for NetApp ONTAP — マルチプロトコルと階層化が強み")
    add_page_number(slide, 5, 9)
    add_footer_bar(slide)

    left_sections = [
        {
            "heading": "✅ メリット",
            "color": COLOR_GREEN,
            "items": [
                "NFS + SMB + iSCSI のマルチプロトコル対応",
                "同一データにNFS/SMB両方からアクセス可能",
                "SSD + Capacity Pool の自動階層化でコスト最適化",
                "ONTAP Snapshot（容量効率が高い）",
                "SnapMirror でクロスリージョン・オンプレ間レプリケーション",
                "データ圧縮・重複排除で実効容量を大幅削減（2〜3倍の効果も）",
                "AWS完全マネージド — インフラ管理不要",
            ]
        },
    ]

    right_sections = [
        {
            "heading": "⚠ デメリット / 注意点",
            "color": COLOR_RED,
            "items": [
                "ONTAP CLIやSystem Managerの操作知識が必要（中程度の学習曲線）",
                "SSD/Capacity Poolの階層化ポリシー設計が必要",
                "スループットキャパシティの費用が大きい（128MB/s: 約$512/月）",
                "Windows管理者にとっては馴染みのない管理体系",
                "費用はFSx Windowsと同等〜やや高い水準",
            ]
        },
    ]

    add_two_column_content(slide, left_sections, right_sections)

    # コスト概算カード
    cost_top = Inches(5.0)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        CONTENT_LEFT, cost_top, Inches(11.8), Inches(1.3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_NAVY
    card.line.color.rgb = COLOR_BORDER
    card.line.width = Pt(0.5)

    add_textbox(
        slide, CONTENT_LEFT + Inches(0.3), cost_top + Inches(0.15),
        Inches(11.2), Inches(0.3),
        "コスト概算（50TB / 東京リージョン / Single-AZ / 階層化構成）",
        size=FONT_BODY, bold=True, color=COLOR_NAVY
    )
    add_textbox(
        slide, CONTENT_LEFT + Inches(0.3), cost_top + Inches(0.5),
        Inches(11.2), Inches(0.6),
        "SSD 5TB: $1,250/月 + Capacity Pool 45TB: $562/月 + スループット(128MB/s): $512/月 → 年間約290〜420万円\n※ 圧縮・重複排除の効果により、実効コストは大幅に下がる可能性あり",
        size=Pt(10), color=COLOR_TEXT
    )

    add_source_note(slide, "出典: AWS公式料金ページ（ap-northeast-1） | ※ 圧縮効率はデータ種別に依存。事前検証推奨")


def slide_06_qumulo(prs):
    """スライド6: Qumulo on AWS 詳細"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "選択肢3: Qumulo on AWS — PB級スケールでは有力だが50TBではオーバースペック")
    add_page_number(slide, 6, 9)
    add_footer_bar(slide)

    left_sections = [
        {
            "heading": "✅ メリット",
            "color": COLOR_GREEN,
            "items": [
                "ペタバイト級まで拡張可能な高いスケーラビリティ",
                "Web UIによるリアルタイムファイルシステム分析・管理",
                "NFS + SMB + S3 API のマルチプロトコル対応",
                "ファイル使用状況・パフォーマンスのリアルタイム可視化",
                "柔軟なカスタマイズが可能",
            ]
        },
    ]

    right_sections = [
        {
            "heading": "⚠ デメリット / 注意点",
            "color": COLOR_RED,
            "items": [
                "AWSマネージドサービスではない（EC2ベースのソフトウェアデプロイ型）",
                "EC2インスタンスの管理（パッチ・OS更新・クラスタ管理）が顧客責任",
                "最低4ノードクラスタが必要 → 50TBスケールではコスト過大",
                "ソフトウェアライセンス + EC2 + EBS の複合課金で高額になりやすい",
                "Qumulo社の Azure版（ANQ）はマネージドだが、AWS版にはマネージド型なし",
                "50TBスケールではFSxの方がコスト効率・運用負荷ともに有利",
            ]
        },
    ]

    add_two_column_content(slide, left_sections, right_sections)

    # 注意カード
    cost_top = Inches(5.0)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        CONTENT_LEFT, cost_top, Inches(11.8), Inches(1.3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFD, 0xF0, 0xE0)  # 薄いオレンジ
    card.line.color.rgb = RGBColor(0xE0, 0xC0, 0x80)
    card.line.width = Pt(0.5)

    add_textbox(
        slide, CONTENT_LEFT + Inches(0.3), cost_top + Inches(0.15),
        Inches(11.2), Inches(0.3),
        "⚠ コスト・提供形態に関する注意",
        size=FONT_BODY, bold=True, color=COLOR_RED
    )
    add_textbox(
        slide, CONTENT_LEFT + Inches(0.3), cost_top + Inches(0.5),
        Inches(11.2), Inches(0.6),
        "AWS上のQumuloはMarketplace経由のEC2デプロイ型であり、FSxのようなマネージドサービスとは本質的に異なります。\n費用はQumulo社またはパートナー経由での個別見積が必要（ライセンス + EC2 + EBS）。50TBスケールではFSxより高額になる可能性が高い",
        size=Pt(10), color=COLOR_TEXT
    )


def slide_07_decision_flow(prs):
    """スライド7: 選択判定フロー"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Windows運用継承 × 50TBスケールなら FSx for Windows が最適解")
    add_page_number(slide, 7, 9)
    add_footer_bar(slide)

    # Q1
    q1_top = Inches(1.5)
    q1_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), q1_top, Inches(7.0), Inches(0.8)
    )
    q1_box.fill.solid()
    q1_box.fill.fore_color.rgb = COLOR_LIGHT_NAVY
    q1_box.line.color.rgb = COLOR_NAVY
    q1_box.line.width = Pt(1)

    add_textbox(
        slide, Inches(1.7), q1_top + Inches(0.1), Inches(6.6), Inches(0.6),
        "Q1: 既存Windowsファイルサーバーの運用方法をそのまま継承したい？ かつSMBのみで十分？",
        size=FONT_BODY, bold=True, color=COLOR_NAVY
    )

    # Q1 YES →
    add_textbox(
        slide, Inches(9.0), q1_top + Inches(0.15), Inches(3.5), Inches(0.5),
        "YES → FSx for Windows ✅",
        size=FONT_BODY_L, bold=True, color=COLOR_GREEN
    )

    # Q2
    q2_top = Inches(2.8)
    add_textbox(
        slide, Inches(3.5), q1_top + Inches(0.8), Inches(2.0), Inches(0.4),
        "NO ↓", size=FONT_BODY, bold=True, color=COLOR_SUBTEXT,
        alignment=PP_ALIGN.CENTER
    )

    q2_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), q2_top, Inches(7.0), Inches(0.9)
    )
    q2_box.fill.solid()
    q2_box.fill.fore_color.rgb = COLOR_LIGHT_NAVY
    q2_box.line.color.rgb = COLOR_NAVY
    q2_box.line.width = Pt(1)

    add_textbox(
        slide, Inches(1.7), q2_top + Inches(0.1), Inches(6.6), Inches(0.7),
        "Q2: NFS + SMBの両方が必要？ または既存NetApp ONTAP環境がある？\n     SnapMirror/Snapshot等のONTAP機能を活用したい？",
        size=FONT_BODY, bold=True, color=COLOR_NAVY
    )

    add_textbox(
        slide, Inches(9.0), q2_top + Inches(0.2), Inches(3.5), Inches(0.5),
        "YES → FSx for ONTAP ✅",
        size=FONT_BODY_L, bold=True, color=COLOR_GREEN
    )

    # Q3
    q3_top = Inches(4.2)
    add_textbox(
        slide, Inches(3.5), q2_top + Inches(0.9), Inches(2.0), Inches(0.4),
        "NO ↓", size=FONT_BODY, bold=True, color=COLOR_SUBTEXT,
        alignment=PP_ALIGN.CENTER
    )

    q3_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), q3_top, Inches(7.0), Inches(0.8)
    )
    q3_box.fill.solid()
    q3_box.fill.fore_color.rgb = COLOR_LIGHT_NAVY
    q3_box.line.color.rgb = COLOR_NAVY
    q3_box.line.width = Pt(1)

    add_textbox(
        slide, Inches(1.7), q3_top + Inches(0.1), Inches(6.6), Inches(0.6),
        "Q3: PB級のスケーラビリティが必要？ リアルタイムファイル分析が必須要件？",
        size=FONT_BODY, bold=True, color=COLOR_NAVY
    )

    add_textbox(
        slide, Inches(9.0), q3_top + Inches(0.15), Inches(3.5), Inches(0.5),
        "YES → Qumulo を検討",
        size=FONT_BODY_L, bold=True, color=RGBColor(0xCC, 0x88, 0x00)
    )

    # 結論
    conclusion_top = Inches(5.5)
    add_textbox(
        slide, Inches(3.5), q3_top + Inches(0.8), Inches(2.0), Inches(0.4),
        "NO ↓", size=FONT_BODY, bold=True, color=COLOR_SUBTEXT,
        alignment=PP_ALIGN.CENTER
    )

    conclusion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), conclusion_top, Inches(10.5), Inches(0.8)
    )
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    conclusion_box.line.color.rgb = COLOR_GREEN
    conclusion_box.line.width = Pt(1.5)

    add_textbox(
        slide, Inches(1.7), conclusion_top + Inches(0.15), Inches(10.1), Inches(0.5),
        "結論: 50TBスケールの場合、FSx for Windows File Server または FSx for NetApp ONTAP で十分対応可能",
        size=FONT_BODY_L, bold=True, color=COLOR_GREEN
    )


def slide_08_cost_comparison(prs):
    """スライド8: コスト比較サマリ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "FSx for Windowsが50TBスケールで最もコスト効率が高い")
    add_page_number(slide, 8, 9)
    add_footer_bar(slide)

    data = [
        ["費用項目", "FSx for Windows\n（HDD / Single-AZ）", "FSx for NetApp ONTAP\n（階層化 / Single-AZ）", "Qumulo on AWS\n（4ノードクラスタ）"],
        ["ストレージ費用\n（月額）", "HDD 50TB:\n$1,250/月", "SSD 5TB: $1,250/月\nCap.Pool 45TB: $562/月", "EBS + ライセンス\n（個別見積）"],
        ["スループット費用\n（月額）", "32MB/s:\n約$176/月", "128MB/s:\n約$512/月", "EC2インスタンス費に含む"],
        ["バックアップ費用", "AWS Backup:\n$0.025/GB/月（増分）", "Snapshot:\nストレージに含む", "独自設計が必要"],
        ["年間概算費用\n（50TB）", "約250〜270万円", "約290〜420万円\n※圧縮効果で変動", "500万円以上の可能性\n（要個別見積）"],
    ]

    col_widths = [Inches(2.2), Inches(3.2), Inches(3.2), Inches(3.2)]
    create_table(slide, data, top=Inches(1.3), height=Inches(4.0), col_widths=col_widths)

    add_source_note(slide, "※ 東京リージョン（ap-northeast-1）、1USD=150円換算。リザーブド価格（1-3年コミット）で最大60%割引可能。最終見積はAWS Pricing Calculator推奨", top=Inches(5.8))

    # 補足
    add_textbox(
        slide, CONTENT_LEFT, Inches(6.2), CONTENT_WIDTH, Inches(0.5),
        "💡 FSx for NetApp ONTAPはデータ圧縮・重複排除により実効コストが大幅に下がる可能性あり。事前にデータ分析を推奨",
        size=Pt(10), bold=False, color=COLOR_ACCENT
    )


def slide_09_next_steps(prs):
    """スライド9: Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "次のステップ")
    add_page_number(slide, 9, 9)
    add_footer_bar(slide)

    # ステップカード
    steps = [
        ("STEP 1", "本資料のご確認", "3つの選択肢のメリット・デメリットをご確認いただき、\nご不明点やご質問があればお知らせください"),
        ("STEP 2", "ヒアリング・要件整理", "お打ち合わせにて以下を詳しくヒアリング:\n• 移行後の運用方針（既存継承 vs 簡素化）\n• プロトコル要件（SMBのみ or NFS併用）\n• バックアップ・DR要件\n• AWS側の体制・リソース確保状況"),
        ("STEP 3", "詳細見積・提案書の作成", "ヒアリング結果をもとに:\n• AWS Pricing Calculator で正確な費用試算\n• 推奨構成の詳細設計\n• 移行計画・スケジュール案の策定"),
        ("STEP 4", "PoC / 検証（必要に応じて）", "選定サービスの検証環境を構築し:\n• パフォーマンステスト\n• 運用手順の事前確認\n• データ移行テスト"),
    ]

    card_top = Inches(1.5)
    for i, (step_num, title, desc) in enumerate(steps):
        y = card_top + Inches(i * 1.4)

        # ステップ番号バッジ
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.75), y, Inches(1.2), Inches(0.4)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_NAVY
        badge.line.fill.background()

        add_textbox(
            slide, Inches(0.75), y + Inches(0.05), Inches(1.2), Inches(0.3),
            step_num, size=Pt(10), bold=True, color=COLOR_WHITE,
            alignment=PP_ALIGN.CENTER
        )

        # タイトル
        add_textbox(
            slide, Inches(2.2), y, Inches(4.0), Inches(0.35),
            title, size=FONT_SUBHEAD, bold=True, color=COLOR_NAVY
        )

        # 説明
        add_textbox(
            slide, Inches(2.2), y + Inches(0.35), Inches(9.5), Inches(0.9),
            desc, size=Pt(10), color=COLOR_TEXT
        )

    # 連絡先案内
    add_textbox(
        slide, CONTENT_LEFT, Inches(7.0), CONTENT_WIDTH, Inches(0.3),
        "ご不明な点がございましたら、お気軽にお問い合わせください。",
        size=FONT_BODY, color=COLOR_SUBTEXT, alignment=PP_ALIGN.CENTER
    )


# ============================================================
# メイン
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 全スライド生成
    slide_01_title(prs)
    slide_02_executive_summary(prs)
    slide_03_comparison_matrix(prs)
    slide_04_fsx_windows(prs)
    slide_05_fsx_ontap(prs)
    slide_06_qumulo(prs)
    slide_07_decision_flow(prs)
    slide_08_cost_comparison(prs)
    slide_09_next_steps(prs)

    # ファイル保存
    if platform.system() == "Linux":
        output_path = "/mnt/c/Users/tshibasaki/Desktop/etc/work/task/00_INBOX/コメダ様_FSx選択肢メリデメ_v2.pptx"
    else:
        output_path = r"C:\Users\tshibasaki\Desktop\etc\work\task\00_INBOX\コメダ様_FSx選択肢メリデメ_v2.pptx"

    output_dir = os.path.dirname(output_path)
    print(f"📁 OS: {platform.system()}")
    print(f"📁 出力ディレクトリ: {output_dir}")
    print(f"📁 ディレクトリ存在: {os.path.exists(output_dir)}")

    try:
        prs.save(output_path)
        saved_size = os.path.getsize(output_path)
        print(f"✅ PPT生成完了: {output_path}")
        print(f"✅ ファイルサイズ: {saved_size} bytes")
        print(f"✅ ファイル存在確認: {os.path.exists(output_path)}")
        print(f"✅ スライド数: {len(prs.slides)}")
    except Exception as e:
        print(f"❌ エラー: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
