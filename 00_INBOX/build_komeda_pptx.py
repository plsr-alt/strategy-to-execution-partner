from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree
import copy, sys
sys.stdout.reconfigure(encoding='utf-8')

src = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310.pptx"
dst = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310_v2.pptx"
prs = Presentation(src)

# ========== ユーティリティ ==========
def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    r_id = sldId.get(qn('r:id'))
    prs.slides.part.drop_rel(r_id)
    sldIdLst.remove(sldId)

def set_para_text(para, text):
    p_elem = para._p
    for r in p_elem.findall(qn('a:r')): p_elem.remove(r)
    for br in p_elem.findall(qn('a:br')): p_elem.remove(br)
    r_elem = etree.SubElement(p_elem, qn('a:r'))
    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text

def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    elem = sldIdLst[from_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(to_idx, elem)

def add_slide_copy(prs, template_idx, insert_at, title_text=None, content_lines=None):
    layout = prs.slide_layouts[2]
    new_slide = prs.slides.add_slide(layout)
    tmpl = prs.slides[template_idx]
    new_sp = new_slide.shapes._spTree
    for child in list(new_sp): new_sp.remove(child)
    for child in tmpl.shapes._spTree: new_sp.append(copy.deepcopy(child))
    move_slide(prs, len(prs.slides)-1, insert_at)
    sl = prs.slides[insert_at]
    if title_text or content_lines:
        for shape in sl.shapes:
            if not shape.has_text_frame: continue
            nm = shape.name
            if title_text and 'title' in nm.lower() or 'タイトル' in nm:
                tf = shape.text_frame; tf.clear()
                p = tf.paragraphs[0]; run = p.add_run()
                run.text = title_text; run.font.size = Pt(18); run.font.bold = True
            elif content_lines and ('content' in nm.lower() or 'コンテンツ' in nm or 'プレースホルダー' in nm):
                tf = shape.text_frame; tf.clear(); tf.word_wrap = True
                first = True
                for text, bold, size, color in content_lines:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    run = p.add_run(); run.text = text
                    run.font.size = Pt(size); run.font.bold = bold
                    if color: run.font.color.rgb = RGBColor(*color)
    return sl

# ===== 削除 =====
for idx in sorted([2,4,5,6,10,11,12,15,16,17,18,19,20,21,22,25], reverse=True):
    delete_slide(prs, idx)
print(f"削除後スライド数: {len(prs.slides)}")

# [0]表紙 [1]はじめに [2]前回ヒアリング [3]Win互換 [4]HA [5]FSxW単価
# [6]費用サマリ [7]利用料概算 [8]今後の進め方(section) [9]確認事項 [10]スケジュール [11]裏表紙

# ===== 表紙 日付修正 =====
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '2025' in para.text and '11' in para.text:
                set_para_text(para, '2026年3月10日')
                print("表紙日付更新")

# ===== 前回ヒアリングスライド(idx=2) 更新 =====
sl3 = prs.slides[2]
for shape in sl3.shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'タイトル' in nm:
        set_para_text(shape.text_frame.paragraphs[0], "前回ご提案内容の振り返り")
    elif 'コンテンツ' in nm or 'プレースホルダー' in nm:
        tf = shape.text_frame; tf.clear(); tf.word_wrap = True
        content = [
            ("前回（2025年11月）のご提案内容を振り返ります。", False, 11, None),
            ("", False, 8, None),
            ("前回ヒアリング内容（主な要件）", True, 13, (0x1F,0x49,0x7D)),
            ("  ・OS：Windows Server", False, 11, None),
            ("  ・ディスク総容量：21TB（使用量：17TB）", False, 11, None),
            ("  ・会話内でのご認識：50TB  ← 本日確認事項", False, 11, (0xC0,0x00,0x00)),
            ("  ・推奨サービス：FSx for Windows File Server", False, 11, None),
            ("", False, 8, None),
            ("前回ご提案した構成（参考）", True, 13, (0x1F,0x49,0x7D)),
            ("  ・FSx for Windows File Server シングルAZ（20TBベース）", False, 11, None),
            ("  ・AWS DataSyncによるデータ移行", False, 11, None),
            ("  ・AWS Backupによるバックアップ", False, 11, None),
            ("  ・NW接続（Direct Connect / VPN）が前提条件", False, 11, (0xC0,0x00,0x00)),
        ]
        first = True
        for text, bold, size, color in content:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run(); run.text = text
            run.font.size = Pt(size); run.font.bold = bold
            if color: run.font.color.rgb = RGBColor(*color)
print("前回ヒアリングスライド更新")

# ===== NW切替支援 セクションヘッダー(idx=3) =====
add_slide_copy(prs, 8, 3, title_text="NW切替支援について")
print("NW切替支援ヘッダー挿入")

# ===== NW切替支援 コンテンツ(idx=4) =====
nw_lines = [
    ("なぜNW切替が必要か", True, 13, (0x1F,0x49,0x7D)),
    ("  FSx for Windows File ServerはAWS VPC内に配置されます。", False, 11, None),
    ("  オンプレミス機器からFSxへアクセスするには AWS接続（NW切替）が前提条件です。", False, 11, None),
    ("", False, 8, None),
    ("接続方式の選択肢", True, 13, (0x1F,0x49,0x7D)),
    ("  Direct Connect（専用線）", False, 11, None),
    ("    安定した帯域・低遅延。大容量データ転送に最適。回線調達に数ヶ月要。", False, 11, None),
    ("  Site-to-Site VPN（インターネットVPN）", False, 11, None),
    ("    低コスト・短期間で構築可能。帯域はインターネット回線に依存。", False, 11, None),
    ("", False, 8, None),
    ("弊社NW切替支援で対応できること", True, 13, (0x1F,0x49,0x7D)),
    ("  ・接続方式の選定支援", False, 11, None),
    ("  ・パラメータシート作成 / 構築 / 疎通確認", False, 11, None),
    ("  ・NW切替完了後にFSx移行支援へシームレスに移行", False, 11, None),
]
add_slide_copy(prs, 2, 4, title_text="NW切替支援の概要", content_lines=nw_lines)
print("NW切替支援コンテンツ挿入")

# 現在の構成:
# [0]表紙 [1]はじめに [2]前回VR [3]NW-header [4]NW-content
# [5]Win互換 [6]HA [7]FSxW単価 [8]費用サマリ [9]利用料概算
# [10]今後の進め方section [11]確認事項 [12]スケジュール [13]裏表紙

# ===== FSxW単価(idx=7) 50TB注記 =====
for shape in prs.slides[7].shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            if 'シングルAZ' in para.text and '参考' in para.text:
                new_p = tf.add_paragraph()
                run = new_p.add_run()
                run.text = "※ 50TB構成の場合、上記単価をもとに別途概算をご提示いたします。"
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
                print("FSxW単価 50TB注記追加")
                break

# ===== 費用サマリ(idx=8) テキスト内20TB→50TB(暫定) =====
for shape in prs.slides[8].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '20TB' in str(run.text):
                    run.text = run.text.replace('20TB', '50TB（暫定）')
                    print("費用サマリ 20TB->50TB(暫定)")
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if '20TB' in str(run.text):
                            run.text = run.text.replace('20TB', '50TB（暫定）')

# ===== 利用料概算(idx=9) 更新 =====
for shape in prs.slides[9].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '20TB' in str(run.text):
                    run.text = run.text.replace('SSD20TB', 'SSD50TB（暫定・確定後に再見積）')
                    print("利用料概算 更新")

# ===== スケジュール(idx=12) 更新 =====
for shape in prs.slides[12].shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'コンテンツ' in nm or 'プレースホルダー' in nm:
        tf = shape.text_frame; tf.clear(); tf.word_wrap = True
        sched = [
            ("今回（3/10）の確認をもとに、以下のスケジュールで進めることを想定しています。", False, 11, None),
            ("弊社でお役に立てる範囲でご協力いたします。", False, 11, None),
            ("", False, 8, None),
            ("【想定スケジュール】", True, 12, (0x1F,0x49,0x7D)),
            ("2026年3月（本日）    NW方式・FSx要件の確認・合意", False, 11, None),
            ("2026年4月～          NW切替支援 着手（詳細設計・構築）", False, 11, None),
            ("2026年5月～          FSx for Windows File Server 設計開始", False, 11, None),
            ("2026年夏頃（目処）   NW切替完了", False, 11, None),
            ("2026年夏～秋         FSx for Windows 環境構築・動作確認", False, 11, None),
            ("2026年秋～冬         データ移行（DataSync）・権限確認", False, 11, None),
            ("2026年末～2027年初   カットオーバー・新環境運用開始", False, 11, None),
            ("2027年3月            既存機器 保守切れ", False, 11, None),
        ]
        first = True
        for text, bold, size, color in sched:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run(); run.text = text
            run.font.size = Pt(size); run.font.bold = bold
            if color: run.font.color.rgb = RGBColor(*color)
        print("スケジュール更新")

# ===== ページ番号更新 =====
for i, sl in enumerate(prs.slides):
    for shape in sl.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t.isdigit() and 1 <= int(t) <= 30 and int(t) != (i+1):
                    set_para_text(para, str(i+1))

prs.save(dst)
print(f"\n保存完了: {dst}")
print(f"総スライド数: {len(prs.slides)}")
for i, sl in enumerate(prs.slides):
    t = ' | '.join([s.text_frame.text[:35].replace('\n',' ').replace('\x0b',' ')
                    for s in sl.shapes if s.has_text_frame and s.text_frame.text.strip()][:2])
    print(f"  [{i+1:02d}] {t}")
