from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree
import sys
sys.stdout.reconfigure(encoding='utf-8')

path = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310_v2.pptx"
prs = Presentation(path)

def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    elem = sldIdLst[from_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(to_idx, elem)

# -----------------------------------------------
# 新スライド追加（末尾 idx=13 として追加 → slide14.xml）
# -----------------------------------------------
layout = prs.slide_layouts[2]  # コンテンツ1
new_sl = prs.slides.add_slide(layout)
print(f"新スライド追加: partname={new_sl.part.partname}")

# -----------------------------------------------
# タイトル設定
# -----------------------------------------------
for shape in new_sl.shapes:
    nm = shape.name
    if not shape.has_text_frame: continue
    if 'Title' in nm or 'タイトル' in nm:
        tf = shape.text_frame; tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = "FSx for Windows File Server 構築ヒアリング"
        run.font.size = Pt(20); run.font.bold = True
        print(f"  タイトル設定: {nm!r}")

# -----------------------------------------------
# 本文設定（6カテゴリ）
# -----------------------------------------------
BLUE  = RGBColor(0x1F, 0x49, 0x7D)
RED   = RGBColor(0xC0, 0x00, 0x00)
BLACK = None

items = [
    # (text, bold, size, color, left_indent)
    # ① ストレージ
    ("① ストレージ容量・構成",                                          True,  12, BLUE,  0),
    ("  ・確認：50TB という認識でよいか（現在17TB使用）",                False, 11, RED,   1),
    ("  ・今後3〜5年の増加見込み",                                       False, 11, BLACK, 1),
    ("  ・ボリューム分割の要否（部門別共有など）",                        False, 11, BLACK, 1),
    ("", False, 6, BLACK, 0),
    # ② Active Directory
    ("② Active Directory 連携",                                         True,  12, BLUE,  0),
    ("  ・ADドメイン名",                                                 False, 11, BLACK, 1),
    ("  ・AD構成：オンプレAD継続利用 ／ AWS Managed AD（どちらか）",      False, 11, RED,   1),
    ("  ・OU / グループポリシーの引き継ぎ要件",                          False, 11, BLACK, 1),
    ("", False, 6, BLACK, 0),
    # ③ アクセス・接続
    ("③ アクセス・接続要件",                                            True,  12, BLUE,  0),
    ("  ・同時接続ユーザー数の目安",                                     False, 11, BLACK, 1),
    ("  ・利用拠点数・リモートアクセスの有無",                           False, 11, BLACK, 1),
    ("", False, 6, BLACK, 0),
    # ④ 可用性
    ("④ 可用性・冗長性",                                                True,  12, BLUE,  0),
    ("  ・シングルAZ ／ マルチAZ（許容できるダウンタイム）",              False, 11, RED,   1),
    ("  ・メンテナンス可能な時間帯",                                     False, 11, BLACK, 1),
    ("", False, 6, BLACK, 0),
    # ⑤ バックアップ
    ("⑤ バックアップ要件",                                              True,  12, BLUE,  0),
    ("  ・バックアップ保持期間・世代数",                                 False, 11, BLACK, 1),
    ("  ・既存バックアップ運用の引き継ぎ要否",                           False, 11, BLACK, 1),
    ("", False, 6, BLACK, 0),
    # ⑥ データ移行
    ("⑥ データ移行",                                                    True,  12, BLUE,  0),
    ("  ・移行元サーバ：台数 ／ OS ／ ファイルシステム",                 False, 11, BLACK, 1),
    ("  ・ACL（アクセス権限）の引き継ぎ要否",                            False, 11, RED,   1),
    ("  ・カットオーバー希望時期・並行運用期間の要件",                   False, 11, BLACK, 1),
]

for shape in new_sl.shapes:
    nm = shape.name
    if not shape.has_text_frame: continue
    if 'Content' in nm or 'コンテンツ' in nm or ('Placeholder' in nm and 'Title' not in nm):
        tf = shape.text_frame; tf.clear(); tf.word_wrap = True
        first = True
        for text, bold, size, color, _ in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if text:
                run = p.add_run()
                run.text = text
                run.font.bold = bold
                run.font.size = Pt(size)
                if color:
                    run.font.color.rgb = color
        print(f"  本文設定: {nm!r}")
        break

# -----------------------------------------------
# 挿入位置を調整: 末尾(13) → idx=6（確認事項の直後）
# -----------------------------------------------
move_slide(prs, 13, 6)

# -----------------------------------------------
# 最終確認
# -----------------------------------------------
print("\n最終スライド構成:")
for i, sl in enumerate(prs.slides):
    txts = [s.text_frame.text[:40].replace(chr(11),' ')
            for s in sl.shapes
            if s.has_text_frame and s.text_frame.text.strip()
            and 'スライド番号' not in s.name]
    t = txts[0][:45] if txts else '(empty)'
    suffix = "  ◀ 補足" if i in [10, 11, 12] else ""
    new_mark = "  ◀ NEW" if i == 6 else ""
    print(f"  [{i+1:02d}] {t!r}{suffix}{new_mark}")

prs.save(path)
print(f"\n上書き保存: {path}")
print(f"総スライド数: {len(prs.slides)}")
