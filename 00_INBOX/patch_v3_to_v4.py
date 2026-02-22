"""
v3 → v4 パッチ
- Slide 16 (idx=15): GUI管理機能スライドを 2カラム版に差し替え
  左: BlueXP（Workload Factory）  右: Lambda リンク
  実際のスクリーンショット付き
- System Manager の記述を除去済み
- 質疑応答は v3 時点で既になし
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SRC  = (r"C:/Users/tshibasaki/Documents/03.案件情報/2025_案件/ウッドワン"
        r"/02-ポストセールス/021-実施ステップ/0211-管理資料"
        r"/AWS-065352-01-報告会資料_20260224_v3.pptx")
DEST = SRC.replace("_v3.pptx", "_v4.pptx")
SS_DIR = r"C:/Users/tshibasaki/Desktop/etc/work/task/00_INBOX/screenshots"

IMG_BLUEXP  = os.path.join(SS_DIR, "screenshot-filesystem-menu-options.png")
IMG_LAMBDA  = os.path.join(SS_DIR, "screenshot-remove-link.png")

prs = Presentation(SRC)
W = prs.slide_width    # 9144000
H = prs.slide_height   # 5143500
print(f"Slide size: {W/914400:.2f}\" x {H/914400:.2f}\"  (v3: {len(prs.slides)} slides)")

# ── ユーティリティ ──────────────────────────────────────

def add_hdr_rect(slide, left, top, width, height, fill_rgb, text, pt=13):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(pt)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

def add_textbox(slide, left, top, width, height, lines):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = para.add_run()
        run.text  = item.get("text", "")
        run.font.size  = Pt(item.get("pt", 11))
        run.font.bold  = item.get("bold", False)
        col = item.get("color")
        if col:
            run.font.color.rgb = col

# ── GUI スライド (idx=15) を差し替え ────────────────────

slide_gui = prs.slides[15]

# タイトル以外のシェイプを全削除
shapes_to_del = []
title_ph = None
for sh in slide_gui.shapes:
    try:
        if sh.placeholder_format and sh.placeholder_format.idx == 0:
            title_ph = sh
            continue
    except Exception:
        pass
    shapes_to_del.append(sh)
for sh in shapes_to_del:
    sh._element.getparent().remove(sh._element)

# タイトルを設定
if title_ph:
    tf = title_ph.text_frame
    if not tf.paragraphs[0].runs:
        tf.paragraphs[0].add_run()
    tf.paragraphs[0].runs[0].text = "FSx for NetApp ONTAP の管理方法"
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(24)

# ── レイアウト定数 (EMU) ──
MARGIN   = Emu(300000)       # 0.33" 左右マージン
GAP      = Emu(180000)       # 0.20" 列間
COL_W    = (W - MARGIN*2 - GAP) // 2   # 1列幅

TITLE_H  = Emu(900000)       # タイトル領域高さ
HDR_H    = Emu(430000)       # カラム見出し高さ
IMG_GAP  = Emu(80000)        # 見出し→画像 間隔

# 画像の縦横比
RATIO_L  = 1372 / 445        # BlueXP  = 3.08
RATIO_R  = 882  / 299        # Lambda  = 2.95

IMG_W    = COL_W
IMG_H_L  = int(IMG_W / RATIO_L)
IMG_H_R  = int(IMG_W / RATIO_R)

DESC_TOP_L = TITLE_H + HDR_H + IMG_GAP + IMG_H_L + Emu(80000)
DESC_TOP_R = TITLE_H + HDR_H + IMG_GAP + IMG_H_R + Emu(80000)
DESC_H     = H - max(DESC_TOP_L, DESC_TOP_R) - Emu(120000)

LEFT_X   = MARGIN
RIGHT_X  = MARGIN + COL_W + GAP

# ═══ 左列: BlueXP (Workload Factory) ═══
BLUE = RGBColor(0x00, 0x5F, 0xB0)

add_hdr_rect(slide_gui, LEFT_X, TITLE_H, COL_W, HDR_H,
             BLUE, "BlueXP  (Workload Factory)")

slide_gui.shapes.add_picture(
    IMG_BLUEXP,
    LEFT_X, TITLE_H + HDR_H + IMG_GAP,
    IMG_W, IMG_H_L
)

add_textbox(slide_gui, LEFT_X, DESC_TOP_L, COL_W, DESC_H, [
    {"text": "NetApp が提供するクラウド管理ポータル", "bold": True,  "pt": 11,
     "color": BLUE},
    {"text": "console.bluexp.netapp.com",             "bold": False, "pt": 10,
     "color": RGBColor(0x44, 0x44, 0x44)},
    {"text": "Console エージェント（NetApp ソフトウェア）をネットワーク内に"
             "配置することで、FSx ファイルシステム・SVM・ボリューム等をGUIで管理可能",
     "bold": False, "pt": 11},
])

# ═══ 右列: Lambda リンク ═══
ORANGE = RGBColor(0xFF, 0x99, 0x00)

add_hdr_rect(slide_gui, RIGHT_X, TITLE_H, COL_W, HDR_H,
             ORANGE, "Lambda リンク（2024〜 新機能）")

slide_gui.shapes.add_picture(
    IMG_LAMBDA,
    RIGHT_X, TITLE_H + HDR_H + IMG_GAP,
    IMG_W, IMG_H_R
)

add_textbox(slide_gui, RIGHT_X, DESC_TOP_R, COL_W, DESC_H, [
    {"text": "Console エージェント不要の新しい接続方式", "bold": True, "pt": 11,
     "color": RGBColor(0xC0, 0x60, 0x00)},
    {"text": "AWS Lambda を使ってリンクを作成。Workload Factory 上から"
             "Console エージェントなしに FSx を管理可能。"
             "AWS Secrets Manager で認証情報を管理できるため"
             "セキュリティ面でも優れている。",
     "bold": False, "pt": 11},
])

print("GUI 2カラムスライド作成完了.")

# ── 保存 & 確認 ──────────────────────────────────────────
prs.save(DEST)
print(f"Saved: {DEST}")

prs2 = Presentation(DEST)
print(f"\n== スライド構成 (v4, {len(prs2.slides)} 枚) ==")
for i, slide in enumerate(prs2.slides, 1):
    t = ""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            t = ph.text_frame.text.strip().replace("\n", "/")[:45]
            break
    if not t:
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip().replace("\n", "/")[:45]
                break
    print(f"  Slide {i:02d}: {t}")
