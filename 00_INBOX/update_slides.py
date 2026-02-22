"""
AWS-065352-01-報告会資料_20260224_v2.pptx 修正スクリプト

変更内容:
1. Slide 02 (本資料の目的) → コンテンツ記入（提案資料ベース）
2. Slide 08 (ご支援結果) 直後 → 「納品物一覧」挿入
3. Slide 14 (今後のご計画について) 直後 → 「FSx for NetApp ONTAP 運用のポイント」挿入
4. さらに後ろ → 「FSx for NetApp ONTAP の GUI 管理機能」挿入
   ・ 箇条書きスタイルは既存スライドに合わせ「●」統一、サブ項目は「・」テキスト
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
import copy
from lxml import etree

PPTX_PATH = (
    r"C:/Users/tshibasaki/Documents/03.案件情報/2025_案件/ウッドワン"
    r"/02-ポストセールス/021-実施ステップ/0211-管理資料"
    r"/AWS-065352-01-報告会資料_20260224_v2.pptx"
)

prs = Presentation(PPTX_PATH)


# ──────────────────────────────────────────────────────
# ユーティリティ
# ──────────────────────────────────────────────────────

def get_title_ph(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph
    return None

def get_body_ph(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    return None

def set_title(slide, text, pt=24):
    ph = get_title_ph(slide)
    if ph is None:
        return
    tf = ph.text_frame
    if not tf.paragraphs[0].runs:
        tf.paragraphs[0].add_run()
    tf.paragraphs[0].runs[0].text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(pt)

def build_body(ph, items):
    """
    items: list of dict
        text  : str
        bold  : bool  (default False)
        pt    : int   (default 13)
        level : int   (default 0)
    既存スライドのスタイルに合わせ、第1パラグラフを再利用して追記。
    """
    tf = ph.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text  = item.get("text", "")
        bold  = item.get("bold", False)
        pt    = item.get("pt", 13)
        level = item.get("level", 0)

        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()

        para.level = level
        if para.runs:
            run = para.runs[0]
        else:
            run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.bold = bold

def move_slide_to(prs, from_idx, to_idx):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide = slides[from_idx]
    xml_slides.remove(slide)
    xml_slides.insert(to_idx, slide)

# コンテンツレイアウト（コンテンツ1 / Slide05 と同じ）を取得
content_layout = prs.slides[4].slide_layout


# ══════════════════════════════════════════════════════
# 1. Slide 02 本資料の目的 にコンテンツを記入
# ══════════════════════════════════════════════════════

slide02 = prs.slides[1]
body_ph = get_body_ph(slide02)

if body_ph is None:
    txBox = slide02.shapes.add_textbox(
        Emu(457200), Emu(1400000), Emu(8229600), Emu(4500000)
    )
    body_ph = txBox

build_body(body_ph, [
    {"text": "本報告会の目的",                                                              "bold": True,  "pt": 14},
    {"text": "本支援（ファイルサーバーの AWS FSx for NetApp ONTAP への移行）の完了をご報告し、今後の安定運用への引継ぎを行うことを目的としています。", "bold": False, "pt": 13, "level": 1},
    {"text": "",                                                                            "pt": 7},
    {"text": "ご支援内容の背景",                                                             "bold": True,  "pt": 14},
    {"text": "既存オンプレファイルサーバーの老朽化・保守期限対応として、AWS FSx for NetApp ONTAP への移行を実施", "bold": False, "pt": 13, "level": 1},
    {"text": "支援範囲：AWS 環境構築 / バックアップ実装（AWS Backup）/ ファイルサーバー移行（AWS DataSync）/ 引継ぎ資料整備", "bold": False, "pt": 13, "level": 1},
    {"text": "",                                                                            "pt": 7},
    {"text": "本日のご報告内容",                                                             "bold": True,  "pt": 14},
    {"text": "背景・目的の確認",       "bold": False, "pt": 13, "level": 1},
    {"text": "ご支援結果の報告",       "bold": False, "pt": 13, "level": 1},
    {"text": "本番切替後の対応事項の共有", "bold": False, "pt": 13, "level": 1},
    {"text": "納品物のご説明",         "bold": False, "pt": 13, "level": 1},
    {"text": "今後のご計画のご提案",   "bold": False, "pt": 13, "level": 1},
])
set_title(slide02, "本資料の目的")
print("✓ Slide 02 updated.")


# ══════════════════════════════════════════════════════
# 2. 「納品物一覧」スライドを Slide 08 の直後に挿入
# ══════════════════════════════════════════════════════

slide_dlv = prs.slides.add_slide(content_layout)
set_title(slide_dlv, "納品物一覧")
body_ph = get_body_ph(slide_dlv)
if body_ph:
    build_body(body_ph, [
        {"text": "以下の資料を納品いたします。", "bold": False, "pt": 13},
        {"text": "",                             "pt": 7},
        {"text": "文書番号",                      "bold": True,  "pt": 13},
        {"text": "AWS-065352-02　AWSパラメータシート",                         "bold": False, "pt": 13, "level": 1},
        {"text": "AWS-065352-03　FSx 運用手順書",                              "bold": False, "pt": 13, "level": 1},
        {"text": "AWS-065352-04　リストア手順書",                              "bold": False, "pt": 13, "level": 1},
        {"text": "AWS-065352-05　移行計画書　（提出済：1/26）",                "bold": False, "pt": 13, "level": 1},
        {"text": "AWS-065352-06　DataSync エージェント導入手順書　（提出済）", "bold": False, "pt": 13, "level": 1},
        {"text": "AWS-065352-07　支援報告書",                                  "bold": False, "pt": 13, "level": 1},
    ])

# Slide 08 = idx 7 → 直後 idx 8 に移動
move_slide_to(prs, len(prs.slides) - 1, 8)
print("✓ 納品物一覧 inserted at index 8.")


# ══════════════════════════════════════════════════════
# 3. 「FSx for NetApp ONTAP 運用のポイント」を
#    今後のご計画 (元 Slide14, 挿入後も idx=14) の直後へ
# ══════════════════════════════════════════════════════

slide_ops = prs.slides.add_slide(content_layout)
set_title(slide_ops, "FSx for NetApp ONTAP 運用のポイント")
body_ph = get_body_ph(slide_ops)
if body_ph:
    build_body(body_ph, [
        {"text": "移行完了後の安定運用のために、以下をご留意ください。", "bold": False, "pt": 13},
        {"text": "",                                                     "pt": 7},
        {"text": "ボリューム容量管理",                                    "bold": True,  "pt": 13},
        {"text": "使用率は 85〜90% 以下を維持することを推奨（90% 超で AWS Backup が失敗する場合あり）", "bold": False, "pt": 12, "level": 1},
        {"text": "定期的なボリューム使用率の監視を実施してください",                                   "bold": False, "pt": 12, "level": 1},
        {"text": "",                                                     "pt": 6},
        {"text": "SSD 拡張は最終手段",                                   "bold": True,  "pt": 13},
        {"text": "SSD 容量の拡張はコスト増大に直結するため、安易な拡張は避けてください",               "bold": False, "pt": 12, "level": 1},
        {"text": "まずデータ整理・削除や、キャパシティプールへのティアリング促進を検討してください",     "bold": False, "pt": 12, "level": 1},
        {"text": "SSD 拡張後は 24 時間程度の安定待機が必要です（ストレージ再配置のため）",              "bold": False, "pt": 12, "level": 1},
        {"text": "",                                                     "pt": 6},
        {"text": "冷却期間（ティアリング）の活用",                        "bold": True,  "pt": 13},
        {"text": "一定期間アクセスのないデータは自動的にキャパシティプール（低コスト HDD 層）へ移動（デフォルト冷却期間：31 日）", "bold": False, "pt": 12, "level": 1},
        {"text": "SSD 拡張前に、冷却期間経過後のデータ移動を待つことを推奨します",                      "bold": False, "pt": 12, "level": 1},
        {"text": "",                                                     "pt": 6},
        {"text": "バックアップ確認",                                     "bold": True,  "pt": 13},
        {"text": "AWS Backup の正常取得・エラー通知を定期確認してください", "bold": False, "pt": 12, "level": 1},
    ])

# 今後のご計画 (元Slide14) は挿入後 idx=14 → 直後 idx=15
move_slide_to(prs, len(prs.slides) - 1, 15)
print("✓ FSx 運用のポイント inserted at index 15.")


# ══════════════════════════════════════════════════════
# 4. 「FSx for NetApp ONTAP の GUI 管理機能」を idx=16 に挿入
#    （運用ポイントスライドの直後）
# ══════════════════════════════════════════════════════

slide_gui = prs.slides.add_slide(content_layout)
set_title(slide_gui, "FSx for NetApp ONTAP の GUI 管理機能")
body_ph = get_body_ph(slide_gui)
if body_ph:
    build_body(body_ph, [
        {"text": "ONTAP System Manager による GUI 操作が可能になりました。", "bold": False, "pt": 13},
        {"text": "",                                                           "pt": 7},
        {"text": "従来の管理方法",                                              "bold": True,  "pt": 13},
        {"text": "BlueXP（現：NetApp Console）経由でのみ ONTAP System Manager にアクセス可能", "bold": False, "pt": 12, "level": 1},
        {"text": "Console エージェント（NetApp ソフトウェア）をネットワーク内に配置する必要があった",     "bold": False, "pt": 12, "level": 1},
        {"text": "",                                                           "pt": 6},
        {"text": "新機能：Lambda リンクの追加（2024〜2025 年）",                "bold": True,  "pt": 13},
        {"text": "AWS Lambda を利用した「リンク」を作成することで、Console エージェントを別途用意せずに ONTAP System Manager を使用可能", "bold": False, "pt": 12, "level": 1},
        {"text": "AWS Secrets Manager で認証情報を管理できるため、BlueXP Workloads に認証情報を保存不要",                               "bold": False, "pt": 12, "level": 1},
        {"text": "NetApp Console（BlueXP）上から同じ UI で操作できるため、学習コストが低い",                                             "bold": False, "pt": 12, "level": 1},
        {"text": "",                                                           "pt": 6},
        {"text": "注意事項",                                                    "bold": True,  "pt": 13},
        {"text": "System Manager を使用するには NetApp Console（BlueXP）へのログインが必要",               "bold": False, "pt": 12, "level": 1},
        {"text": "System Manager での変更は AWS コンソールへの反映に数分かかる場合あり",                  "bold": False, "pt": 12, "level": 1},
        {"text": "第 2 世代ファイルシステム（複数 HA ペア）では NetApp Console は非対応",                  "bold": False, "pt": 12, "level": 1},
    ])

# FSx運用ポイントは idx=15 → GUI管理は idx=16
move_slide_to(prs, len(prs.slides) - 1, 16)
print("✓ FSx GUI 管理スライド inserted at index 16.")


# ══════════════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════════════
SAVE_PATH = PPTX_PATH.replace("_v2.pptx", "_v3.pptx")
prs.save(SAVE_PATH)
print("\nSaved: " + SAVE_PATH)

# 確認用: スライド一覧
prs2 = Presentation(PPTX_PATH)
print("\n== スライド構成 ==")
for i, slide in enumerate(prs2.slides, 1):
    title = ""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            title = ph.text_frame.text.strip().replace("\n", "/")[:40]
            break
    if not title:
        for shape in slide.shapes:
            if shape.has_text_frame:
                title = shape.text_frame.text.strip().replace("\n", "/")[:40]
                break
    print(f"  Slide {i:02d}: {title}")
