"""
ウッドワン報告会資料 v2 - 本文フォントサイズを 14pt に統一
対象: コンテンツスライドの本文のみ（中表紙・表紙・裏表紙は除外）
"""
from pptx import Presentation
from pptx.util import Pt

SRC = (
    r"C:/Users/tshibasaki/Documents/03.案件情報/2025_案件/ウッドワン"
    r"/02-ポストセールス/021-実施ステップ/0211-管理資料"
    r"/AWS-065352-01-報告会資料_20260224_v2.pptx"
)
DST = SRC  # 上書き保存

TARGET_SIZE = Pt(14)

# スキップするスライド（0-indexed）
# 0=表紙, 1=本資料の目的(中表紙), 2=アジェンダ(中表紙), 3=背景振り返り(中表紙),
# 6=ご報告(中表紙), 12=今後のご計画(中表紙), 14=質疑応答(中表紙), 15=裏表紙
SKIP_SLIDES = {0, 1, 2, 3, 6, 12, 14, 15}

# スキップする placeholder type（数値）
# 1=TITLE, 3=CENTER_TITLE, 13=SLIDE_NUMBER
SKIP_PH_TYPES = {1, 3, 13}

prs = Presentation(SRC)
changed = 0

for slide_idx, slide in enumerate(prs.slides):

    # ===== 質疑応答スライドの見出しを設定 =====
    if slide_idx == 14:
        for shape in slide.shapes:
            if shape.name == 'Text Placeholder 1' and shape.has_text_frame:
                tf = shape.text_frame
                tf.text = '質疑応答'
                run = tf.paragraphs[0].runs[0]
                run.font.size = Pt(32)
                run.font.bold = True
        continue  # フォントサイズ変更はスキップ

    # 中表紙・表紙・裏表紙はスキップ
    if slide_idx in SKIP_SLIDES:
        continue

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # タイトル・スライド番号はスキップ
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in SKIP_PH_TYPES:
                continue

        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                # 20pt以上（中表紙の大見出し等）はスキップ
                if run.font.size and run.font.size >= Pt(20):
                    continue
                if run.font.size != TARGET_SIZE:
                    run.font.size = TARGET_SIZE
                    changed += 1

print(f"変更した run 数: {changed}")

prs.save(DST)
print("保存完了:", DST)
