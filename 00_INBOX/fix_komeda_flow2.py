from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from lxml import etree
import sys
sys.stdout.reconfigure(encoding='utf-8')

path = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310_v2.pptx"
prs = Presentation(path)

def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    elem = sldIdLst[from_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(to_idx, elem)

def clear_and_set(tf, lines):
    tf.clear(); tf.word_wrap = True
    first = True
    for text, bold, size in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if text:
            run = p.add_run()
            run.text = text
            run.font.bold = bold
            run.font.size = Pt(size)

# -----------------------------------------------
# 現構成（fix_komeda_flow.py 実行後）
# [0]表紙 [1]はじめに [2]前回VR [3]Win互換 [4]HA
# [5]NW-header [6]NW-content [7]FSxW単価
# [8]費用サマリ [9]利用料概算 [10]確認事項
# [11]スケジュール [12]裏表紙
# -----------------------------------------------
# 目標構成
# [0]表紙 [1]はじめに [2]前回VR
# [3]NW-header [4]NW-content [5]確認事項
# [6]費用サマリ [7]利用料概算 [8]スケジュール
# ---- 補足 ----
# [9]Win互換 [10]HA [11]FSxW単価
# [12]裏表紙
# -----------------------------------------------

# Step1: NW-header [5] → [3]
move_slide(prs, 5, 3)
# [0]表紙 [1]はじめに [2]前回VR [3]NW-header [4]Win互換 [5]HA [6]NW-content ...

# Step2: NW-content [6] → [4]
move_slide(prs, 6, 4)
# [0]表紙 [1]はじめに [2]前回VR [3]NW-header [4]NW-content [5]Win互換 [6]HA [7]FSxW単価 [8]費用サマリ [9]利用料概算 [10]確認事項 [11]スケジュール [12]裏表紙

# Step3: 確認事項 [10] → [5]
move_slide(prs, 10, 5)
# [5]確認事項 [6]Win互換 [7]HA [8]FSxW単価 [9]費用サマリ [10]利用料概算 [11]スケジュール [12]裏表紙

# Step4: 費用サマリ [9] → [6]
move_slide(prs, 9, 6)
# [6]費用サマリ [7]Win互換 [8]HA [9]FSxW単価 [10]利用料概算 [11]スケジュール [12]裏表紙

# Step5: 利用料概算 [10] → [7]
move_slide(prs, 10, 7)
# [7]利用料概算 [8]Win互換 [9]HA [10]FSxW単価 [11]スケジュール [12]裏表紙

# Step6: スケジュール [11] → [8]
move_slide(prs, 11, 8)
# [8]スケジュール [9]Win互換 [10]HA [11]FSxW単価 [12]裏表紙


# -----------------------------------------------
# はじめに（idx=1）: ヒアリング目的に合わせたアジェンダへ
# -----------------------------------------------
for shape in prs.slides[1].shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'コンテンツ' in nm or 'Content' in nm or ('プレースホルダー' in nm and 'タイトル' not in nm and 'スライド' not in nm):
        if 'アジェンダ' in shape.text_frame.text or 'NW切替' in shape.text_frame.text or '振り返り' in shape.text_frame.text:
            clear_and_set(shape.text_frame, [
                ("前回ご提案の内容を踏まえ、本日は受注に向けた要件確認・ヒアリングを目的としております。", False, 11),
                ("", False, 9),
                ("■ 本日のアジェンダ", True, 12),
                ("　１．前回ご提案内容の振り返り", False, 11),
                ("　２．NW切替支援について（移行の前提条件）", False, 11),
                ("　３．今回の確認・ヒアリング事項", False, 11),
                ("　４．費用概算・今後のスケジュール確認", False, 11),
                ("", False, 9),
                ("■ 本日のゴール", True, 12),
                ("　・NW切替の進め方・接続方式について認識合わせする", False, 11),
                ("　・FSxW見積確定に必要な要件（容量・構成等）をヒアリングする", False, 11),
                ("　・2026年5月のFSx設計開始に向けた合意を得る", False, 11),
            ])
            print("はじめに アジェンダ更新（ヒアリング重視）")
            break

# -----------------------------------------------
# NW-content（idx=4）: 橋渡し文をヒアリングへの導線に更新
# -----------------------------------------------
for shape in prs.slides[4].shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'Content' in nm or 'コンテンツ' in nm or ('Placeholder' in nm and 'Title' not in nm):
        tf = shape.text_frame
        if '接続方式' in tf.text or 'NW切替が必要か' in tf.text:
            # 末尾の橋渡し文を更新（前回追加分を置き換え）
            paras = tf.paragraphs
            # 末尾3段落を確認して「NW切替完了後のステップ」ブロックを削除・更新
            full = tf.text
            if 'NW切替完了後のステップ' in full:
                # 末尾の追記パラグラフを除去してから再追記
                p_elem = tf._txBody
                all_p = p_elem.findall(qn('a:p'))
                # 「NW切替完了後のステップ」以降を削除
                remove_flag = False
                for p in all_p:
                    txt = ''.join(r.text or '' for r in p.findall('.//' + qn('a:t')))
                    if 'NW切替完了後のステップ' in txt:
                        remove_flag = True
                    if remove_flag:
                        p_elem.remove(p)
            # 新しい橋渡し文を追加
            p = tf.add_paragraph()
            p_elem2 = p._p
            r = etree.SubElement(p_elem2, qn('a:r'))
            etree.SubElement(r, qn('a:t'))  # 空行
            for text, bold, size in [
                ("", False, 9),
                ("以上を踏まえ、次のページにて本日確認させていただきたい事項をご説明します。", False, 11),
            ]:
                p = tf.add_paragraph()
                if text:
                    run = p.add_run()
                    run.text = text
                    run.font.bold = bold
                    run.font.size = Pt(size)
            print("NW-content 橋渡し文 → ヒアリングへの導線に更新")
            break

# -----------------------------------------------
# 最終確認
# -----------------------------------------------
print("\n最終スライド構成:")
for i, sl in enumerate(prs.slides):
    txts = [s.text_frame.text[:40].replace(chr(11),' ')
            for s in sl.shapes
            if s.has_text_frame and s.text_frame.text.strip() and 'スライド番号' not in s.name]
    t = txts[0][:45] if txts else '(empty)'
    suffix = "  ◀ 補足" if i in [9, 10, 11] else ""
    print(f"  [{i+1:02d}] {t!r}{suffix}")

prs.save(path)
print(f"\n上書き保存: {path}")
print(f"総スライド数: {len(prs.slides)}")
