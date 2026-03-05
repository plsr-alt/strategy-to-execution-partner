#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
コメダ様 ファイルサーバー移行提案資料 — テンプレート形式対応版 v4
方針: プレースホルダーをXMLから完全削除 → テキストボックスで全コンテンツを配置
これによりテンプレートのゴーストテキストとの重なりを完全に排除
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import os
import platform


# ============================================================
# テンプレート読み込み
# ============================================================

def load_template():
    if platform.system() == "Linux":
        path = "/mnt/c/Users/tshibasaki/Desktop/etc/work/task/00_INBOX/template_komeda_20260225.pptx"
    else:
        path = r"C:\Users\tshibasaki\Desktop\etc\work\task\00_INBOX\template_komeda_20260225.pptx"
    return Presentation(path)


# ============================================================
# コア: プレースホルダー完全削除
# ============================================================

def remove_all_placeholders(slide):
    """スライド上の全プレースホルダーをXMLから削除する。
    スライドマスターの背景装飾は残るが、ゴーストテキストは完全に消える。"""
    spTree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        sp = ph._element
        spTree.remove(sp)


def add_slide_clean(prs, layout_index):
    """レイアウトからスライドを追加し、全プレースホルダーを削除して返す"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    remove_all_placeholders(slide)
    return slide


# ============================================================
# ヘルパー: シェイプ追加
# ============================================================

def _set_ea_font(run, typeface='メイリオ'):
    """東アジアフォントを設定"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', typeface)


def add_textbox(slide, left_in, top_in, width_in, height_in, text,
                font_size=Pt(11), bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = "Meiryo"
    if color:
        run.font.color.rgb = color
    _set_ea_font(run)
    return txBox


def add_title(slide, text):
    """タイトルテキストボックス（テンプレートのタイトル位置に配置）"""
    add_textbox(slide, 0.31, 0.20, 8.07, 0.55, text,
                font_size=Pt(18), bold=True)


def add_page_number(slide, number):
    """ページ番号テキストボックス"""
    add_textbox(slide, 9.33, 5.28, 0.60, 0.30, str(number),
                font_size=Pt(11), alignment=PP_ALIGN.CENTER)


def add_bullets(slide, left_in, top_in, width_in, items, font_size=Pt(11)):
    """箇条書きテキストボックスを追加"""
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(4.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = font_size
        run.font.name = "Meiryo"
        _set_ea_font(run)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    return txBox


# --- テーブル ---

def _set_cell_color(cell, color):
    tcPr = cell._tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))


def _set_cell_border(cell, border_pos, color, width=Pt(0.5)):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    border_map = {'top': 'a:lnT', 'bottom': 'a:lnB', 'left': 'a:lnL', 'right': 'a:lnR'}
    tag = qn(border_map[border_pos])
    for old in tcPr.findall(tag):
        tcPr.remove(old)
    ln = etree.SubElement(tcPr, tag)
    ln.set('w', str(int(width)))
    solidFill = etree.SubElement(ln, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))


def create_table(slide, data, left=Inches(0.43), top=Inches(0.87),
                 width=Inches(9.2), height=Inches(4.2)):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    col_w = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_w

    header_color = RGBColor(0x1F, 0x49, 0x7D)
    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = data[0][j]
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Meiryo"
        _set_cell_color(cell, header_color)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = data[i][j]
            run.font.size = Pt(9)
            run.font.name = "Meiryo"
            bg = RGBColor(0xF2, 0xF2, 0xF2) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            _set_cell_color(cell, bg)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            for pos in ['top', 'bottom', 'left', 'right']:
                _set_cell_border(cell, pos, RGBColor(0xD9, 0xD9, 0xD9))

    return table


# ============================================================
# スライド生成
# ============================================================

def slide_01_title(prs, page_num):
    """表紙"""
    slide = add_slide_clean(prs, 0)  # 表紙レイアウト

    add_textbox(slide, 4.65, 2.50, 5.35, 0.80,
                "AWS ファイルストレージ選択肢\nメリット・デメリット比較",
                font_size=Pt(18), bold=True)

    add_textbox(slide, 4.65, 3.40, 5.35, 0.35,
                "FSx for Windows / FSx for NetApp ONTAP / Qumulo",
                font_size=Pt(9))

    add_textbox(slide, 4.49, 4.78, 5.28, 0.67,
                "2026年3月\n株式会社アシスト\nビジネスインフラ技術本部",
                font_size=Pt(10))

    add_textbox(slide, 0.31, 2.20, 4.0, 0.40,
                "株式会社コメダ 御中", font_size=Pt(14), bold=True)

    return page_num + 1


def slide_02_toc(prs, page_num):
    """アジェンダ"""
    slide = add_slide_clean(prs, 1)  # コンテンツレイアウト
    add_title(slide, "アジェンダ")
    add_page_number(slide, page_num)

    add_bullets(slide, 0.43, 0.87, 9.0, [
        "背景・現状",
        "3つの選択肢の概要と特性",
        "総合比較マトリクス",
        "各選択肢の詳細（メリット・デメリット）",
        "コスト比較",
        "選択判定フロー",
        "Next Steps",
    ])
    return page_num + 1


def slide_03_background(prs, page_num):
    """背景・現状"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "背景・現状")
    add_page_number(slide, page_num)

    add_bullets(slide, 0.43, 0.87, 9.0, [
        "現行: オンプレミス Windows ファイルサーバー（約50TB）",
        "課題: インフラ老朽化、運用体制の効率化が必要",
        "方針: AWSクラウド環境へのファイルサーバー移行を検討中",
        "重要な判断基準:",
        "  - 既存運用体制との継続性（Windowsスキル継承）",
        "  - 総所有コスト（3年間の費用）",
        "  - 運用保守の手間（学習曲線）",
        "  - スケーラビリティと柔軟性",
    ])
    return page_num + 1


def slide_04_overview(prs, page_num):
    """3選択肢の概要"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "3つのAWSファイルストレージサービス")
    add_page_number(slide, page_num)

    create_table(slide, [
        ["選択肢", "サービス種別", "対応プロトコル", "特徴"],
        ["1. FSx for\nWindows FS", "完全マネージド\n（AWS純正）", "SMB", "既存Windows運用を継承\n管理が簡単"],
        ["2. FSx for\nNetApp ONTAP", "完全マネージド\n（AWS純正）", "NFS/SMB\niSCSI", "マルチプロトコル\n階層化でコスト最適化"],
        ["3. Qumulo\non AWS", "セルフマネージド\n（EC2ベース）", "NFS/SMB\nS3 API", "PB級スケール\nリアルタイム分析機能"],
    ], top=Inches(0.87), height=Inches(3.5))
    return page_num + 1


def slide_05_comparison_matrix(prs, page_num):
    """総合比較マトリクス"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "総合比較マトリクス")
    add_page_number(slide, page_num)

    create_table(slide, [
        ["評価項目", "FSx Windows", "FSx ONTAP", "Qumulo"],
        ["運用難易度\n（Windows継承）", "◎ 低", "○ 中", "△ 中〜高"],
        ["マネージド度", "◎ 完全MG", "◎ 完全MG", "× セルフMG"],
        ["初期導入コスト\n（50TB）", "中", "中", "高（クラスタ必須）"],
        ["年間ランニング\nコスト（50TB）", "約250〜270万", "約290〜420万", "要個別見積"],
        ["スケーラビリティ", "最大512TB\n(HDD)", "無制限\n(Capacity Pool)", "PB級"],
        ["推奨用途", "Windows既存継承\nSMB専用", "NFS必須\nマルチプロトコル", "超大規模\nPB級スケール"],
    ], top=Inches(0.87), height=Inches(4.2))
    return page_num + 1


def slide_06_fsx_windows(prs, page_num):
    """FSx for Windows 詳細"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "選択肢1: FSx for Windows File Server")
    add_page_number(slide, page_num)

    add_textbox(slide, 0.43, 0.87, 4.3, 0.30, "✅ メリット", font_size=Pt(12), bold=True)
    add_bullets(slide, 0.43, 1.20, 4.3, [
        "既存Windowsファイルサーバー運用をそのまま継承",
        "MMC / PowerShell など馴染みのツールで操作",
        "Active Directory完全統合（細かいACL制御）",
        "AWS Backupで自動日次バックアップ",
        "AWS完全マネージド — インフラ管理不要",
    ], font_size=Pt(10))

    add_textbox(slide, 5.20, 0.87, 4.3, 0.30, "⚠ デメリット", font_size=Pt(12), bold=True)
    add_bullets(slide, 5.20, 1.20, 4.3, [
        "SMB のみ（NFS非対応）",
        "SSD構成は高額 → HDD推奨",
        "リストアは新規FSとして復元",
        "費用: 年間約250〜270万円",
        "50TBスケールなら費用効率が高い",
    ], font_size=Pt(10))
    return page_num + 1


def slide_07_fsx_ontap(prs, page_num):
    """FSx for NetApp ONTAP 詳細"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "選択肢2: FSx for NetApp ONTAP")
    add_page_number(slide, page_num)

    add_textbox(slide, 0.43, 0.87, 4.3, 0.30, "✅ メリット", font_size=Pt(12), bold=True)
    add_bullets(slide, 0.43, 1.20, 4.3, [
        "NFS + SMB + iSCSI のマルチプロトコル対応",
        "同一データにNFS/SMB両方からアクセス可能",
        "SSD + Capacity Pool で自動階層化・コスト最適化",
        "ONTAP Snapshot（容量効率が高い）",
        "SnapMirrorでレプリケーション対応",
        "圧縮・重複排除で実効容量を2〜3倍削減可能",
    ], font_size=Pt(10))

    add_textbox(slide, 5.20, 0.87, 4.3, 0.30, "⚠ デメリット", font_size=Pt(12), bold=True)
    add_bullets(slide, 5.20, 1.20, 4.3, [
        "ONTAP CLIの操作知識が必要",
        "階層化ポリシー設計が複雑",
        "スループット費用が大きい",
        "Windows管理者に馴染みない",
        "費用: 年間約290〜420万円",
    ], font_size=Pt(10))
    return page_num + 1


def slide_08_qumulo(prs, page_num):
    """Qumulo on AWS 詳細"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "選択肢3: Qumulo on AWS")
    add_page_number(slide, page_num)

    add_textbox(slide, 0.43, 0.87, 4.3, 0.30, "✅ メリット", font_size=Pt(12), bold=True)
    add_bullets(slide, 0.43, 1.20, 4.3, [
        "PB級まで拡張可能な高いスケーラビリティ",
        "Webの管理UI — リアルタイム分析機能",
        "NFS + SMB + S3 API のマルチプロトコル",
        "ファイル使用状況・パフォーマンスを可視化",
        "柔軟なカスタマイズが可能",
    ], font_size=Pt(10))

    add_textbox(slide, 5.20, 0.87, 4.3, 0.30, "⚠ デメリット", font_size=Pt(12), bold=True)
    add_bullets(slide, 5.20, 1.20, 4.3, [
        "AWS マネージドサービスではない（EC2ベース）",
        "EC2管理が顧客責任（パッチ・更新等）",
        "最低4ノードクラスタ必須 → 高額",
        "50TBではオーバースペック",
        "費用: 要個別見積（高額になりやすい）",
    ], font_size=Pt(10))
    return page_num + 1


def slide_09_cost_comparison(prs, page_num):
    """コスト比較"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "コスト比較（50TB / 東京リージョン）")
    add_page_number(slide, page_num)

    create_table(slide, [
        ["費用項目", "FSx Windows\n(HDD/Single-AZ)", "FSx ONTAP\n(SSD+Cap.Pool)"],
        ["ストレージ\n(月額)", "$1,250\n(HDD 50TB)", "SSD: $1,250\nCap.Pool: $562"],
        ["スループット\n(月額)", "$176\n(32MB/s)", "$512\n(128MB/s)"],
        ["バックアップ\n(月額)", "$1,250\n(増分課金)", "Snapshotに含む"],
        ["年間合計\n(概算)", "約250〜270万円", "約290〜420万円"],
        ["学習コスト", "低（Windows継承）", "中（ONTAP技術習得）"],
        ["推奨", "✅ 50TBスケール\nではこれ！", "NFS必須時の\n第一選択"],
    ], top=Inches(0.87), height=Inches(4.2))
    return page_num + 1


def slide_10_decision_flow(prs, page_num):
    """選択判定フロー"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "選択判定フロー")
    add_page_number(slide, page_num)

    add_textbox(slide, 0.43, 0.95, 6.5, 0.35,
                "Q1: 既存Windows運用を継承？ SMBのみで十分？",
                font_size=Pt(11), bold=True)
    add_textbox(slide, 7.20, 0.95, 2.5, 0.35,
                "YES → FSx for Windows ✅", font_size=Pt(10), bold=True)

    add_textbox(slide, 2.50, 1.35, 1.5, 0.25, "NO ↓", font_size=Pt(9))

    add_textbox(slide, 0.43, 1.75, 6.5, 0.35,
                "Q2: NFS + SMBの両方が必要？ 既存NetApp環境がある？",
                font_size=Pt(11), bold=True)
    add_textbox(slide, 7.20, 1.75, 2.5, 0.35,
                "YES → FSx for ONTAP ✅", font_size=Pt(10), bold=True)

    add_textbox(slide, 2.50, 2.15, 1.5, 0.25, "NO ↓", font_size=Pt(9))

    add_textbox(slide, 0.43, 2.55, 6.5, 0.35,
                "Q3: PB級のスケーラビリティが必須要件？",
                font_size=Pt(11), bold=True)
    add_textbox(slide, 7.20, 2.55, 2.5, 0.35,
                "YES → Qumulo を検討", font_size=Pt(10), bold=True)

    add_textbox(slide, 0.43, 3.30, 9.0, 0.50,
                "結論: 50TBスケール × Windows継承 → FSx for Windows が最適解",
                font_size=Pt(13), bold=True)
    return page_num + 1


def slide_11_next_steps(prs, page_num):
    """Next Steps"""
    slide = add_slide_clean(prs, 1)
    add_title(slide, "Next Steps")
    add_page_number(slide, page_num)

    add_bullets(slide, 0.43, 0.87, 9.0, [
        "STEP 1: 本資料の確認・ご質問",
        "STEP 2: 詳細ヒアリング（運用方針、プロトコル要件、バックアップ要件等）",
        "STEP 3: 正確な見積・詳細設計（AWS Pricing Calculator活用）",
        "STEP 4: PoC/検証環境の構築（パフォーマンステスト、運用手順確認）",
        "STEP 5: 移行計画・スケジュール決定",
    ])
    return page_num + 1


def slide_12_back_cover(prs):
    """裏表紙"""
    slide = add_slide_clean(prs, 4)  # 裏表紙レイアウト

    add_textbox(slide, 3.0, 1.5, 4.0, 0.4,
                "お問合わせ先", font_size=Pt(18), bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 3.0, 2.2, 4.0, 0.4,
                "株式会社アシスト", font_size=Pt(14), bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 3.0, 2.7, 4.0, 0.3,
                "ビジネスインフラ技術本部", font_size=Pt(11), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 3.0, 3.2, 4.0, 0.3,
                "URL  https://www.ashisuto.co.jp/", font_size=Pt(10), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.3, 9.0, 0.6,
                "記載されている価格は、別段の記載がない限り税抜価格です。別途消費税がかかります。\n"
                "記載されている製品およびサービスは、提供している各社の商標または登録商標です。",
                font_size=Pt(7), alignment=PP_ALIGN.CENTER)


# ============================================================
# メイン
# ============================================================

def main():
    prs = load_template()

    # テンプレートの既存スライドを全削除（マスター・レイアウトは保持）
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # スライド生成
    pn = 1
    pn = slide_01_title(prs, pn)
    pn = slide_02_toc(prs, pn)
    pn = slide_03_background(prs, pn)
    pn = slide_04_overview(prs, pn)
    pn = slide_05_comparison_matrix(prs, pn)
    pn = slide_06_fsx_windows(prs, pn)
    pn = slide_07_fsx_ontap(prs, pn)
    pn = slide_08_qumulo(prs, pn)
    pn = slide_09_cost_comparison(prs, pn)
    pn = slide_10_decision_flow(prs, pn)
    pn = slide_11_next_steps(prs, pn)
    slide_12_back_cover(prs)

    # ファイル保存
    if platform.system() == "Linux":
        output_path = "/mnt/c/Users/tshibasaki/Documents/03.案件情報/2026/コメダ/064195_ファイルサーバ構築/01-プリセールス/012-提案資料/【コメダ様】ファイルサーバー移行提案資料_FSx比較版.pptx"
    else:
        output_path = r"C:\Users\tshibasaki\Documents\03.案件情報\2026\コメダ\064195_ファイルサーバ構築\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバー移行提案資料_FSx比較版.pptx"

    print(f"📁 OS: {platform.system()}")
    print(f"📁 出力: {output_path}")

    try:
        prs.save(output_path)
        saved_size = os.path.getsize(output_path)
        print(f"✅ PPT生成完了!")
        print(f"✅ ファイルサイズ: {saved_size} bytes")
        print(f"✅ スライド数: {len(prs.slides)}")
    except Exception as e:
        print(f"❌ エラー: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
