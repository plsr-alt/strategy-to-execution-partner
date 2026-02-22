from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree
import sys, re
sys.stdout.reconfigure(encoding='utf-8')

SRC = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310.pptx"
DST = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310_v2.pptx"
prs = Presentation(SRC)

BLUE  = RGBColor(0x1F, 0x49, 0x7D)
RED   = RGBColor(0xC0, 0x00, 0x00)

# ===== ユーティリティ =====
def delete_slide(prs, index):
    lst = prs.slides._sldIdLst
    sid = lst[index]
    prs.slides.part.drop_rel(sid.get(qn('r:id')))
    lst.remove(sid)

def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    e = lst[from_idx]; lst.remove(e); lst.insert(to_idx, e)

def clear_para(para):
    p = para._p
    for tag in ['a:r','a:br','a:fld']:
        for e in p.findall(qn(tag)): p.remove(e)

def set_para_text(para, text):
    clear_para(para)
    if not text: return
    r = etree.SubElement(para._p, qn('a:r'))
    t = etree.SubElement(r, qn('a:t'))
    t.text = text

def set_run(tf, lines, clear=True):
    if clear: tf.clear(); tf.word_wrap = True
    first = True
    for text, bold, size, color in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if text:
            run = p.add_run(); run.text = text
            run.font.bold = bold; run.font.size = Pt(size)
            if color: run.font.color.rgb = color

def find_shape(slide, *keys):
    for s in slide.shapes:
        nm = s.name
        if any(k in nm for k in keys) and s.has_text_frame:
            return s
    return None

# ===== STEP1: 新スライドを先に追加（28枚時 → slide29/30/31、衝突なし） =====
nw_header  = prs.slides.add_slide(prs.slide_layouts[1])  # slide29
nw_content = prs.slides.add_slide(prs.slide_layouts[2])  # slide30
fsx_hearing= prs.slides.add_slide(prs.slide_layouts[2])  # slide31
print(f"NW-header:   {nw_header.part.partname}")
print(f"NW-content:  {nw_content.part.partname}")
print(f"FSx-hearing: {fsx_hearing.part.partname}")

# NW-header テキスト（堀越さん担当のセクションヘッダー）
for s in nw_header.shapes:
    if s.has_text_frame:
        tf = s.text_frame; tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = "NW切替支援について"
        run.font.bold = True; run.font.size = Pt(28)
        break

# NW-content（堀越さん担当 - プレースホルダーとして記載）
for s in nw_content.shapes:
    if not s.has_text_frame: continue
    nm = s.name
    if 'Title' in nm or 'タイトル' in nm:
        tf = s.text_frame; tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = "NW切替支援の概要"
        run.font.size = Pt(18); run.font.bold = True
    elif 'Content' in nm or 'コンテンツ' in nm or 'Placeholder' in nm:
        set_run(s.text_frame, [
            ("なぜNW切替が必要か",                                           True,  13, BLUE),
            ("  FSx for Windows File ServerはAWS VPC内に配置されます。",      False, 11, None),
            ("  オンプレミスからFSxへのアクセスにはAWS接続（NW切替）が前提です。", False, 11, None),
            ("", False, 7, None),
            ("接続方式の選択肢",                                              True,  13, BLUE),
            ("  Direct Connect（専用線）― 安定帯域・低遅延。回線調達に数ヶ月要。", False, 11, None),
            ("  Site-to-Site VPN（インターネットVPN）― 低コスト・短期構築可能。",  False, 11, None),
            ("", False, 7, None),
            ("弊社NW切替支援の範囲",                                          True,  13, BLUE),
            ("  ・接続方式の選定支援",                                         False, 11, None),
            ("  ・パラメータシート作成 / 構築 / 疎通確認",                      False, 11, None),
            ("  ・NW切替完了後、FSx移行支援へシームレスに移行",                  False, 11, None),
            ("", False, 7, None),
            ("以上を踏まえ、次のページにて本日確認させていただきたい事項をご説明します。", False, 11, None),
        ])

# FSx-hearing スライド内容設定（削除前に実施）
for s in fsx_hearing.shapes:
    if not s.has_text_frame: continue
    nm = s.name
    if 'Title' in nm or 'タイトル' in nm:
        tf = s.text_frame; tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = "FSx for Windows File Server 構築ヒアリング"
        run.font.size = Pt(20); run.font.bold = True
    elif 'Content' in nm or 'コンテンツ' in nm or ('Placeholder' in nm and 'Title' not in nm):
        set_run(s.text_frame, [
            ("① ストレージ容量・構成",                                       True,  12, BLUE),
            ("  ・確認：50TB という認識でよいか（現在17TB使用）",               False, 11, RED),
            ("  ・今後3〜5年の増加見込み",                                    False, 11, None),
            ("  ・ボリューム分割の要否（部門別共有など）",                      False, 11, None),
            ("", False, 6, None),
            ("② Active Directory 連携",                                     True,  12, BLUE),
            ("  ・ADドメイン名",                                              False, 11, None),
            ("  ・AD構成：オンプレAD継続利用 ／ AWS Managed AD",               False, 11, RED),
            ("  ・OU / グループポリシーの引き継ぎ要件",                        False, 11, None),
            ("", False, 6, None),
            ("③ アクセス・接続要件",                                          True,  12, BLUE),
            ("  ・同時接続ユーザー数の目安",                                   False, 11, None),
            ("  ・利用拠点数・リモートアクセスの有無",                          False, 11, None),
            ("", False, 6, None),
            ("④ 可用性・冗長性",                                              True,  12, BLUE),
            ("  ・シングルAZ ／ マルチAZ（許容できるダウンタイム）",             False, 11, RED),
            ("  ・メンテナンス可能な時間帯",                                   False, 11, None),
            ("", False, 6, None),
            ("⑤ バックアップ要件",                                            True,  12, BLUE),
            ("  ・保持期間・世代数",                                           False, 11, None),
            ("  ・既存バックアップ運用の引き継ぎ要否",                          False, 11, None),
            ("", False, 6, None),
            ("⑥ データ移行",                                                  True,  12, BLUE),
            ("  ・移行元サーバ：台数 ／ OS ／ ファイルシステム",               False, 11, None),
            ("  ・ACL（アクセス権限）の引き継ぎ要否",                           False, 11, RED),
            ("  ・カットオーバー希望時期・並行運用期間",                        False, 11, None),
        ])

# ===== STEP2: 16スライド削除（元28枚のインデックス） =====
to_delete = sorted([2,4,5,6,10,11,12,15,16,17,18,19,20,21,22,25], reverse=True)
for idx in to_delete:
    delete_slide(prs, idx)
print(f"削除後: {len(prs.slides)} slides")
# 残存(0-based): [0]表紙 [1]はじめに [2]前回VR [3]Win互換 [4]HA [5]FSxW単価
#               [6]費用サマリ [7]利用料概算 [8]今後section [9]確認事項
#               [10]スケジュール [11]裏表紙 [12]NW-header [13]NW-content [14]FSx-hearing

# ===== STEP3: 表紙 日付 =====
for s in prs.slides[0].shapes:
    if s.has_text_frame:
        for p in s.text_frame.paragraphs:
            if '2025' in p.text and '11月' in p.text:
                set_para_text(p, '2026年3月10日')

# ===== STEP4: はじめに（ヒアリング重視アジェンダ） =====
for s in prs.slides[1].shapes:
    if not s.has_text_frame: continue
    nm = s.name
    if 'コンテンツ' in nm or 'Content' in nm or ('プレースホルダー' in nm and 'タイトル' not in nm and 'スライド' not in nm):
        if 'アジェンダ' in s.text_frame.text or 'NW切替' in s.text_frame.text or '振り返り' in s.text_frame.text:
            set_run(s.text_frame, [
                ("前回ご提案の内容を踏まえ、本日は受注に向けた要件確認・ヒアリングを目的としております。", False, 11, None),
                ("", False, 8, None),
                ("■ 本日のアジェンダ",                        True,  12, None),
                ("　１．前回ご提案内容の振り返り",              False, 11, None),
                ("　２．NW切替支援について（移行の前提条件）",   False, 11, None),
                ("　３．今回の確認・ヒアリング事項",             False, 11, None),
                ("　４．費用概算・今後のスケジュール確認",       False, 11, None),
                ("", False, 8, None),
                ("■ 本日のゴール",                            True,  12, None),
                ("　・NW切替の進め方・接続方式について認識合わせする",              False, 11, None),
                ("　・FSxW見積確定に必要な要件（容量・構成等）をヒアリングする",    False, 11, None),
                ("　・2026年5月のFSx設計開始に向けた合意を得る",                  False, 11, None),
            ])
            break

# ===== STEP5: 前回VR =====
sl = prs.slides[2]
for s in sl.shapes:
    if not s.has_text_frame: continue
    nm = s.name
    if 'タイトル' in nm or 'Title' in nm:
        set_para_text(s.text_frame.paragraphs[0], "前回ご提案内容の振り返り")
    elif 'コンテンツ' in nm or 'Content' in nm or ('プレースホルダー' in nm and 'タイトル' not in nm and 'スライド' not in nm):
        set_run(s.text_frame, [
            ("前回（2025年11月）のご提案内容を振り返ります。", False, 11, None),
            ("", False, 7, None),
            ("前回ヒアリング内容（主な要件）",              True,  13, BLUE),
            ("  OS：Windows Server",                    False, 11, None),
            ("  ディスク総容量：21TB（使用量：17TB）",       False, 11, None),
            ("  会話内でのご認識：50TB ← 本日確認事項",      False, 11, RED),
            ("  推奨サービス：FSx for Windows File Server", False, 11, None),
            ("", False, 7, None),
            ("前回ご提案した構成（参考）",                  True,  13, BLUE),
            ("  FSx for Windows File Server シングルAZ（20TBベース）", False, 11, None),
            ("  AWS DataSyncによるデータ移行",              False, 11, None),
            ("  AWS Backupによるバックアップ",              False, 11, None),
            ("  NW接続（Direct Connect / VPN）が前提条件",  False, 11, RED),
        ])

# ===== STEP6: FSxW単価 50TB注記 =====
for s in prs.slides[5].shapes:
    if s.has_text_frame:
        for p in s.text_frame.paragraphs:
            if 'シングルAZ' in p.text and '参考' in p.text:
                new_p = s.text_frame.add_paragraph()
                run = new_p.add_run()
                run.text = "※ 50TB構成の場合、上記単価をもとに別途概算をご提示いたします。"
                run.font.size = Pt(11); run.font.color.rgb = RED
                break

# ===== STEP7: 費用サマリ 20TB→50TB =====
for s in prs.slides[6].shapes:
    if s.has_text_frame:
        for p in s.text_frame.paragraphs:
            for r in p.runs:
                if '20TB' in str(r.text): r.text = r.text.replace('20TB','50TB（暫定）')
    if s.has_table:
        for row in s.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        if '20TB' in str(r.text): r.text = r.text.replace('20TB','50TB（暫定）')

# ===== STEP8: 利用料概算 =====
for s in prs.slides[7].shapes:
    if s.has_text_frame:
        for p in s.text_frame.paragraphs:
            for r in p.runs:
                if 'SSD20TB' in str(r.text):
                    r.text = r.text.replace('SSD20TB','SSD50TB（暫定・確定後に再見積）')

# ===== STEP9: スケジュール =====
for s in prs.slides[10].shapes:
    if not s.has_text_frame: continue
    nm = s.name
    if 'コンテンツ' in nm or 'Content' in nm or ('プレースホルダー' in nm and 'タイトル' not in nm and 'スライド' not in nm):
        set_run(s.text_frame, [
            ("今回（3/10）の確認をもとに、以下スケジュールで進めることを想定しています。", False, 11, None),
            ("", False, 7, None),
            ("【想定スケジュール】",                                       True,  12, BLUE),
            ("2026年3月（本日）    NW方式・FSx要件の確認・合意",            False, 11, None),
            ("2026年4月〜          NW切替支援 着手（詳細設計・構築）",      False, 11, None),
            ("2026年5月〜          FSx for Windows File Server 設計開始", False, 11, None),
            ("2026年夏頃（目処）   NW切替完了",                             False, 11, None),
            ("2026年夏〜秋         FSx for Windows 環境構築・動作確認",     False, 11, None),
            ("2026年秋〜冬         データ移行（DataSync）・権限確認",        False, 11, None),
            ("2026年末〜2027年初   カットオーバー・新環境運用開始",          False, 11, None),
            ("2027年3月            既存機器 保守切れ",                      False, 11, None),
        ])

# ===== STEP10: スライド順序調整 =====
# 現在(0-based): [0]表紙[1]はじめに[2]前回VR[3]Win互換[4]HA[5]FSxW単価
#               [6]費用サマリ[7]利用料概算[8]今後section[9]確認事項
#               [10]スケジュール[11]裏表紙[12]NW-header[13]NW-content
# 目標: NW-header→3, NW-content→4, 今後section削除, 確認→5後へ

# 今後section(idx=8)削除
delete_slide(prs, 8)
# 現在: [8]確認事項 [9]スケジュール [10]裏表紙 [11]NW-header [12]NW-content

# NW-header[11]→3, NW-content[12]→4
move_slide(prs, 11, 3)
move_slide(prs, 12, 4)
# 現在: [0]表紙[1]はじめに[2]前回VR[3]NW-header[4]NW-content
#       [5]Win互換[6]HA[7]FSxW単価[8]費用サマリ[9]利用料概算
#       [10]確認事項[11]スケジュール[12]裏表紙

# 確認事項[10]→5（NW-contentの直後へ）
move_slide(prs, 10, 5)
# 費用サマリ[9]→6
move_slide(prs, 9, 6)
# 利用料概算[10]→7
move_slide(prs, 10, 7)
# スケジュール[11]→8
move_slide(prs, 11, 8)
# 現在: [0]表紙[1]はじめに[2]前回VR[3]NW-header[4]NW-content
#       [5]確認事項[6]費用サマリ[7]利用料概算[8]スケジュール
#       [9]Win互換[10]HA[11]FSxW単価[12]裏表紙 = 13枚

# ===== STEP11: FSx-hearingを確認事項の直後(idx=6)へ移動 =====
# STEP10終了後: [0]表紙...[5]確認事項[6]費用サマリ[7]利用料概算[8]スケジュール
#               [9]Win互換[10]HA[11]FSxW単価[12]裏表紙[13]FSx-hearing
move_slide(prs, 13, 6)
# 最終: [0]表紙[1]はじめに[2]前回VR[3]NW-header[4]NW-content
#       [5]確認事項[6]FSxヒアリング[7]費用サマリ[8]利用料概算[9]スケジュール
#       [10]Win互換[11]HA[12]FSxW単価[13]裏表紙 = 14枚

# ===== 最終確認 =====
print(f"\n最終スライド構成（{len(prs.slides)}枚）:")
for i, sl in enumerate(prs.slides):
    txts = [s.text_frame.text[:35].replace(chr(11),' ')
            for s in sl.shapes
            if s.has_text_frame and s.text_frame.text.strip()
            and 'スライド番号' not in s.name]
    t = txts[0][:40] if txts else '(empty)'
    mark = " ◀ 補足" if i in [10,11,12] else ""
    print(f"  [{i+1:02d}] {t!r}{mark}")

prs.save(DST)
print(f"\n保存完了: {DST}")
