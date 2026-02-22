from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy, sys
sys.stdout.reconfigure(encoding='utf-8')

SRC = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\ウッドワン\02-ポストセールス\021-実施ステップ\0211-管理資料\AWS-065352-01-報告会資料_20260224.pptx"
DST = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\ウッドワン\02-ポストセールス\021-実施ステップ\0211-管理資料\AWS-065352-01-報告会資料_20260224_v2.pptx"

prs = Presentation(SRC)

# ===== ユーティリティ =====
def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    e = lst[from_idx]; lst.remove(e); lst.insert(to_idx, e)

def reorder_slides(prs, new_order):
    """new_order: 現在のインデックスを目標順に並べたリスト"""
    lst = prs.slides._sldIdLst
    elements = [lst[i] for i in new_order]
    for e in elements:
        lst.remove(e)
    for e in elements:
        lst.append(e)

def find_shape(slide, *name_keys):
    for s in slide.shapes:
        if any(k in s.name for k in name_keys) and s.has_text_frame:
            return s
    return None

def set_tf(tf, lines):
    """lines: [(text, bold, size)] """
    tf.clear(); tf.word_wrap = True
    first = True
    for text, bold, size in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if text:
            r = p.add_run(); r.text = text
            r.font.bold = bold; r.font.size = Pt(size)

def set_para(para, text):
    p = para._p
    for tag in ['a:r','a:br','a:fld']:
        for e in p.findall(qn(tag)): p.remove(e)
    if text:
        r = etree.SubElement(p, qn('a:r'))
        t = etree.SubElement(r, qn('a:t'))
        t.text = text

def replace_google_shape_text(slide, old_text_fragment, new_text):
    """Google Shape 系（表紙など）のテキスト置換"""
    for s in slide.shapes:
        if s.has_text_frame and old_text_fragment in s.text_frame.text:
            tf = s.text_frame
            tf.clear(); tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = new_text
            return True
    return False

# ===== STEP1: 新スライドを先に追加（partname 衝突防止） =====
# 14: 本資料の目的 (コンテンツ1 = layout index 1)
# 15: 中表紙 質疑応答 (中表紙－パターン1 = layout index 4)
slide_mokuteki  = prs.slides.add_slide(prs.slide_layouts[1])   # idx 14
slide_qanda_hdr = prs.slides.add_slide(prs.slide_layouts[4])   # idx 15
print(f"本資料の目的:  {slide_mokuteki.part.partname}")
print(f"質疑応答ヘッダ: {slide_qanda_hdr.part.partname}")

# ----- 本資料の目的 -----
for s in slide_mokuteki.shapes:
    if not s.has_text_frame: continue
    if 'タイトル' in s.name or 'Title' in s.name:
        set_tf(s.text_frame, [("本資料の目的", True, 20)])
    elif 'プレースホルダー' in s.name and 'タイトル' not in s.name and 'スライド' not in s.name:
        set_tf(s.text_frame, [
            ("この度は、ファイルサーバーのAWS移行に伴う環境構築・移行支援をご発注いただき誠にありがとうございます。", False, 12),
            ("", False, 8),
            ("本資料は、2025年11月〜2026年1月にわたる支援内容の完了報告および", False, 12),
            ("納品資料のご説明を目的としております。", False, 12),
            ("", False, 8),
            ("■ 本日のアジェンダ", True, 12),
            ("　1. 背景・目的の振り返り", False, 11),
            ("　2. ご報告（支援結果・対応事項・納品資料）", False, 11),
            ("　3. 今後のご計画について", False, 11),
            ("　4. 質疑応答", False, 11),
        ])

# ----- 質疑応答ヘッダー -----
for s in slide_qanda_hdr.shapes:
    if s.has_text_frame and 'テキスト プレースホルダー' in s.name:
        set_tf(s.text_frame, [("質疑応答", True, 32)])
        break

# ===== STEP2: 既存スライドのテキスト更新 =====

# --- slide 0: 表紙 ---
sl = prs.slides[0]
replace_google_shape_text(sl, "現状報告/本番切替について", "ファイルサーバー移行\n支援完了報告")
replace_google_shape_text(sl, "2025年1月22日", "2026年2月24日\n株式会社アシスト\nビジネスインフラ技術本部")

# --- slide 1: アジェンダ ---
sl = prs.slides[1]
for s in sl.shapes:
    if not s.has_text_frame: continue
    if "差分移行" in s.text_frame.text or "本番切替日" in s.text_frame.text:
        set_tf(s.text_frame, [
            ("1. 背景・目的の振り返り", False, 12),
            ("2. ご報告", False, 12),
            ("　　・ご支援結果", False, 11),
            ("　　・本番切替後の対応事項", False, 11),
            ("　　・納品資料一覧", False, 11),
            ("3. 今後のご計画について", False, 12),
            ("4. 質疑応答", False, 12),
        ])

# --- slide 2: 中表紙「背景・目的の振り返り」---
sl = prs.slides[2]
for s in sl.shapes:
    if s.has_text_frame and "差分移行" in s.text_frame.text:
        set_tf(s.text_frame, [("背景・目的の振り返り", True, 32)])

# --- slide 3: 背景・目的①（背景・前提）---
sl = prs.slides[3]
for s in sl.shapes:
    if not s.has_text_frame: continue
    if s.name == 'タイトル 1':
        set_tf(s.text_frame, [("ご依頼背景・目的の認識合わせ（背景・前提）", True, 18)])
    elif s.name == 'テキスト プレースホルダー 3':
        set_tf(s.text_frame, [
            ("本支援ご依頼に至る背景", True, 13),
            ("  ・既存オンプレファイルサーバーの老朽化・保守期限への対応が課題", False, 11),
            ("  ・AWS FSx for NetApp ONTAP への安全かつ効率的な移行が目的", False, 11),
            ("", False, 7),
            ("移行対象・前提", True, 13),
            ("  ・移行先サービス：Amazon FSx for NetApp ONTAP", False, 11),
            ("  ・移行ツール：AWS DataSync", False, 11),
            ("  ・移行元データ量：上限 13TB", False, 11),
            ("  ・前提：移行元環境と移行先 VPC がプライベート通信可能であること", False, 11),
            ("", False, 7),
            ("弊社（株式会社アシスト）支援範囲", True, 13),
            ("  ・AWS 環境構築支援（FSx for NetApp ONTAP 構築）", False, 11),
            ("  ・バックアップ実装支援（AWS Backup）", False, 11),
            ("  ・ファイルサーバー移行支援（DataSync による差分同期〜本番切替）", False, 11),
            ("  ・引継ぎ資料整備（運用手順書・リストア手順書 等）", False, 11),
        ])
# 表が残っていたら非表示にするため削除
for s in sl.shapes:
    if s.name == '表 4':
        sp = s._element
        sp.getparent().remove(sp)

# --- slide 4: 背景・目的②（ご支援範囲・ご期待）---
sl = prs.slides[4]
for s in sl.shapes:
    if not s.has_text_frame: continue
    if s.name == 'タイトル 1':
        set_tf(s.text_frame, [("ご依頼背景・目的の認識合わせ（ご期待範囲）", True, 18)])
    elif s.name == 'テキスト プレースホルダー 3':
        set_tf(s.text_frame, [
            ("① AWS 環境構築支援", True, 13),
            ("  FSx for NetApp ONTAP ファイルシステム・SVM・ボリューム・共有・クォータ設定", False, 11),
            ("", False, 7),
            ("② バックアップ実装支援", True, 13),
            ("  AWS Backup によるバックアップ設定・テスト・リストア検証", False, 11),
            ("", False, 7),
            ("③ ファイルサーバー移行支援", True, 13),
            ("  AWS DataSync を用いた差分同期（12月〜1月）および本番切替（1/24）", False, 11),
            ("  ACL・フォルダ構成の整合性を維持した安全なデータ移行", False, 11),
            ("", False, 7),
            ("④ 引継ぎ資料整備", True, 13),
            ("  運用手順書・リストア手順書・パラメータシート・移行計画書 等", False, 11),
        ])
for s in sl.shapes:
    if s.name == '表 6':
        sp = s._element
        sp.getparent().remove(sp)

# --- slide 5（旧：エラー対応 BackUp_on）→ ご支援結果 ---
sl = prs.slides[5]
for s in sl.shapes:
    if not s.has_text_frame: continue
    if s.name == 'タイトル 1':
        set_tf(s.text_frame, [("ご支援結果", True, 20)])
    elif s.name == 'テキスト プレースホルダー 3':
        set_tf(s.text_frame, [
            ("各技術支援の実施結果は以下のとおりです。", False, 11),
            ("", False, 8),
            ("【AWS 環境構築支援】　　完了 ✓", True, 12),
            ("  ・FSx for NetApp ONTAP 構築（2025年11月）", False, 11),
            ("  ・Audit ボリューム・クォータ設定", False, 11),
            ("", False, 7),
            ("【バックアップ実装支援】　完了 ✓", True, 12),
            ("  ・AWS Backup によるバックアップ設定・リストア検証（2025年11月）", False, 11),
            ("", False, 7),
            ("【ファイルサーバー移行支援】　完了 ✓", True, 12),
            ("  ・DataSync 差分同期（2025年12月〜2026年1月）", False, 11),
            ("  ・本番切替（DNS切替）：2026年1月24日 完了", False, 11),
        ])
for s in sl.shapes:
    if s.name == '表 6':
        sp = s._element
        sp.getparent().remove(sp)

# --- slide 6（旧：エラー残対応）→ 納品資料一覧 ---
sl = prs.slides[6]
# テキストボックス・表 等を削除してクリーン化
to_remove = []
for s in sl.shapes:
    if s.name in ['表 6', 'テキスト ボックス 4']:
        to_remove.append(s._element)
for e in to_remove:
    e.getparent().remove(e)
for s in sl.shapes:
    if not s.has_text_frame: continue
    if s.name == 'タイトル 1':
        set_tf(s.text_frame, [("納品資料一覧", True, 20)])
    elif s.name == 'テキスト プレースホルダー 3':
        set_tf(s.text_frame, [
            ("以下の資料を納品いたします。", False, 11),
            ("", False, 8),
            ("文書番号 / ドキュメント名", True, 12),
            ("  AWS-065352-02　AWSパラメータシート", False, 11),
            ("  AWS-065352-03　FSx 運用手順書", False, 11),
            ("  AWS-065352-04　リストア手順書", False, 11),
            ("  AWS-065352-05　移行計画書　　　　　　　　　　　　　　（提出済：1/26）", False, 11),
            ("  AWS-065352-06　DataSync エージェント導入手順書　　（提出済）", False, 11),
            ("  AWS-065352-07　支援報告書", False, 11),
        ])

# --- slide 7（中表紙 本番切替）→ 中表紙「ご報告」---
sl = prs.slides[7]
for s in sl.shapes:
    if s.has_text_frame and "本番切替" in s.text_frame.text:
        set_tf(s.text_frame, [("ご報告", True, 32)])

# --- slide 8（タイムスケジュール）→ 今後のご計画 ---
sl = prs.slides[8]
for s in sl.shapes:
    if not s.has_text_frame: continue
    if s.name == 'タイトル 2':
        set_tf(s.text_frame, [("今後のご計画について", True, 20)])
    elif s.name == 'テキスト プレースホルダー 3':
        set_tf(s.text_frame, [
            ("本支援の完了後、以下についてご検討・ご対応をお願いいたします。", False, 11),
            ("", False, 8),
            ("【お客様側で実施いただくこと】", True, 12),
            ("  ・運用手順書・リストア手順書の内容確認および社内周知", False, 11),
            ("  ・AWS Backup の正常取得・エラー通知の定期確認", False, 11),
            ("  ・ボリューム使用率の定期監視（目安：85〜90%以下を維持）", False, 11),
            ("", False, 8),
            ("【ご不明点・追加のご支援が必要な場合】", True, 12),
            ("  お気軽にお問い合わせください。", False, 11),
            ("  woodone_fs_prj@ashisuto.co.jp", False, 11),
        ])

# --- slide 9（中表紙 その他）→ 中表紙「今後のご計画について」---
sl = prs.slides[9]
for s in sl.shapes:
    if s.has_text_frame and "その他" in s.text_frame.text:
        set_tf(s.text_frame, [("今後のご計画について", True, 32)])

# --- slide 10: クォータ → タイトルのみ微修正 ---
sl = prs.slides[10]
for s in sl.shapes:
    if s.name == 'タイトル 2' and 'クォータ' in s.text_frame.text:
        set_tf(s.text_frame, [("本番切替後の対応事項①：クォータ", True, 18)])

# --- slide 11: ボリューム → タイトル微修正 ---
sl = prs.slides[11]
for s in sl.shapes:
    if s.name == 'タイトル 1' and 'ボリューム' in s.text_frame.text:
        set_tf(s.text_frame, [("本番切替後の対応事項②：ボリュームサイズ変更", True, 18)])

# --- slide 12: AWS Backup → タイトル微修正 ---
sl = prs.slides[12]
for s in sl.shapes:
    if s.name == 'タイトル 1' and 'Backup' in s.text_frame.text:
        set_tf(s.text_frame, [("本番切替後の対応事項③：AWS Backup", True, 18)])

# ===== STEP3: スライド順序を目標順に並び替え =====
# 現在の indices: 0〜15（新スライド 14=本資料の目的, 15=質疑応答ヘッダー）
# 目標: [0,14, 1, 2, 3, 4, 7, 5, 10, 11, 12, 6, 9, 8, 15, 13]
# [01]表紙, [02]本資料の目的, [03]アジェンダ, [04]背景ヘッダー, [05]背景①, [06]背景②,
# [07]ご報告ヘッダー, [08]ご支援結果, [09]クォータ, [10]ボリューム, [11]Backup,
# [12]納品資料, [13]今後のヘッダー, [14]今後の計画, [15]質疑応答ヘッダー, [16]裏表紙
reorder_slides(prs, [0, 14, 1, 2, 3, 4, 7, 5, 10, 11, 12, 6, 9, 8, 15, 13])

# ===== 最終確認 =====
print(f"\n最終スライド構成（{len(prs.slides)}枚）:")
for i, sl in enumerate(prs.slides):
    txts = [s.text_frame.text[:40].replace('\n',' ').replace('\x0b',' ')
            for s in sl.shapes
            if s.has_text_frame and s.text_frame.text.strip()
            and 'スライド番号' not in s.name]
    t = txts[0] if txts else '(empty)'
    print(f"  [{i+1:02d}] {t!r}")

prs.save(DST)
print(f"\n保存完了: {DST}")
