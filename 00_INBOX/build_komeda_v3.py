from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree
import sys
sys.stdout.reconfigure(encoding='utf-8')

src = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310.pptx"
dst = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310_v2.pptx"
prs = Presentation(src)

# ========== ユーティリティ ==========
def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    r_id = sldId.get(qn('r:id'))
    prs.slides.part.drop_rel(r_id)
    sldIdLst.remove(sldId)

def set_para_text(para, text):
    p_elem = para._p
    for r in p_elem.findall(qn('a:r')): p_elem.remove(r)
    for br in p_elem.findall(qn('a:br')): p_elem.remove(br)
    for fld in p_elem.findall(qn('a:fld')): p_elem.remove(fld)
    if not text: return
    r_elem = etree.SubElement(p_elem, qn('a:r'))
    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text

def set_run(tf, lines, clear=True):
    if clear:
        tf.clear()
        tf.word_wrap = True
    first = True
    for text, bold, size, color in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if not text:
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)

def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    elem = sldIdLst[from_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(to_idx, elem)

# ===== STEP1: 新スライドを先に追加（28枚時に追加→slide29/30で衝突なし）=====
# layout[1]=中表紙, layout[2]=コンテンツ
layout_sec = prs.slide_layouts[1]
nw_header_slide = prs.slides.add_slide(layout_sec)   # idx=28 → slide29.xml
layout_con = prs.slide_layouts[2]
nw_content_slide = prs.slides.add_slide(layout_con)  # idx=29 → slide30.xml
print(f"新スライド追加後: {len(prs.slides)} slides")
print(f"  NW-header partname: {nw_header_slide.part.partname}")
print(f"  NW-content partname: {nw_content_slide.part.partname}")

# ===== STEP2: NW-header タイトル設定 =====
for shape in nw_header_slide.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "NW切替支援について"
        run.font.bold = True
        run.font.size = Pt(28)
        print(f"  NW-header shape set: {shape.name!r} → {run.text!r}")
        break

# ===== STEP3: NW-content 設定 =====
for shape in nw_content_slide.shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'タイトル' in nm or 'title' in nm.lower() or 'Title' in nm:
        tf = shape.text_frame; tf.clear()
        p = tf.paragraphs[0]; run = p.add_run()
        run.text = "NW切替支援の概要"
        run.font.size = Pt(18); run.font.bold = True
        print(f"  NW-content title: {run.text!r}")
    elif 'コンテンツ' in nm or 'Content' in nm or ('Placeholder' in nm and 'Title' not in nm):
        set_run(shape.text_frame, [
            ("なぜNW切替が必要か", True, 13, (0x1F,0x49,0x7D)),
            ("  FSx for Windows File ServerはAWS VPC内に配置されます。", False, 11, None),
            ("  オンプレミス機器からFSxへアクセスするにはAWS接続（NW切替）が前提条件です。", False, 11, None),
            ("", False, 8, None),
            ("接続方式の選択肢", True, 13, (0x1F,0x49,0x7D)),
            ("  Direct Connect（専用線）", False, 11, None),
            ("    安定した帯域・低遅延。大容量データ転送に最適。回線調達に数ヶ月要。", False, 11, None),
            ("  Site-to-Site VPN（インターネットVPN）", False, 11, None),
            ("    低コスト・短期間で構築可能。帯域はインターネット回線に依存。", False, 11, None),
            ("", False, 8, None),
            ("弊社NW切替支援で対応できること", True, 13, (0x1F,0x49,0x7D)),
            ("  接続方式の選定支援", False, 11, None),
            ("  パラメータシート作成 / 構築 / 疎通確認", False, 11, None),
            ("  NW切替完了後にFSx移行支援へシームレスに移行", False, 11, None),
        ])
        print(f"  NW-content body set: {shape.name!r}")

# ===== STEP4: 削除（後ろから）元28枚のインデックスを対象 =====
# 削除対象（0-based, 元の28枚のうち）:
#   2=課題・背景, 4=FS紹介, 5=AWSストレージ, 6=ストレージ,
#   10=FSxN主な特徴, 11=FSxN単価, 12=技術支援AWS利用料,
#   15=構成図, 16=前提条件, 17=FS環境構築, 18-20=空白,
#   21=FS移行支援, 22=空白, 25=今後の検討ポイント
to_delete = sorted([2,4,5,6,10,11,12,15,16,17,18,19,20,21,22,25], reverse=True)
for idx in to_delete:
    delete_slide(prs, idx)
print(f"削除後: {len(prs.slides)} slides")
# 現構成(0-based): [0]表紙 [1]はじめに [2]前回VR [3]Win互換 [4]HA [5]FSxW単価
#                  [6]費用サマリ [7]利用料概算 [8]今後section [9]確認事項
#                  [10]スケジュール [11]裏表紙 [12]NW-header [13]NW-content

# ===== STEP5: 表紙 日付 =====
for shape in prs.slides[0].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if '2025' in para.text and '11月' in para.text:
                set_para_text(para, '2026年3月10日')
print("表紙日付更新")

# ===== STEP6: 前回ヒアリング(idx=2) =====
sl = prs.slides[2]
for shape in sl.shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'タイトル' in nm or 'Title' in nm:
        set_para_text(shape.text_frame.paragraphs[0], "前回ご提案内容の振り返り")
    elif 'コンテンツ' in nm or 'Content' in nm or ('Placeholder' in nm and 'Title' not in nm and 'スライド' not in nm):
        set_run(shape.text_frame, [
            ("前回（2025年11月）のご提案内容を振り返ります。", False, 11, None),
            ("", False, 8, None),
            ("前回ヒアリング内容（主な要件）", True, 13, (0x1F,0x49,0x7D)),
            ("  OS：Windows Server", False, 11, None),
            ("  ディスク総容量：21TB（使用量：17TB）", False, 11, None),
            ("  会話内でのご認識：50TB　← 本日確認事項", False, 11, (0xC0,0x00,0x00)),
            ("  推奨サービス：FSx for Windows File Server", False, 11, None),
            ("", False, 8, None),
            ("前回ご提案した構成（参考）", True, 13, (0x1F,0x49,0x7D)),
            ("  FSx for Windows File Server シングルAZ（20TBベース）", False, 11, None),
            ("  AWS DataSyncによるデータ移行", False, 11, None),
            ("  AWS Backupによるバックアップ", False, 11, None),
            ("  NW接続（Direct Connect / VPN）が前提条件", False, 11, (0xC0,0x00,0x00)),
        ])
print("前回VRスライド更新")

# ===== STEP7: FSxW単価(idx=5) 50TB注記 =====
for shape in prs.slides[5].shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            if 'シングルAZ' in para.text and '参考' in para.text:
                new_p = tf.add_paragraph()
                run = new_p.add_run()
                run.text = "※ 50TB構成の場合、上記単価をもとに別途概算をご提示いたします。"
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
                break
print("FSxW単価 50TB注記追加")

# ===== STEP8: 費用サマリ(idx=6) 20TB→50TB(暫定) =====
for shape in prs.slides[6].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '20TB' in str(run.text):
                    run.text = run.text.replace('20TB', '50TB（暫定）')
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if '20TB' in str(run.text):
                            run.text = run.text.replace('20TB', '50TB（暫定）')
print("費用サマリ 20TB→50TB更新")

# ===== STEP9: 利用料概算(idx=7) =====
for shape in prs.slides[7].shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'SSD20TB' in str(run.text):
                    run.text = run.text.replace('SSD20TB', 'SSD50TB（暫定・確定後に再見積）')
print("利用料概算 更新")

# ===== STEP10: スケジュール(idx=10) =====
for shape in prs.slides[10].shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'コンテンツ' in nm or 'Content' in nm or ('Placeholder' in nm and 'Title' not in nm and 'スライド' not in nm):
        set_run(shape.text_frame, [
            ("今回（3/10）の確認をもとに、以下スケジュールで進めることを想定しています。", False, 11, None),
            ("弊社でお役に立てる範囲でご協力いたします。", False, 11, None),
            ("", False, 8, None),
            ("【想定スケジュール】", True, 12, (0x1F,0x49,0x7D)),
            ("2026年3月（本日）    NW方式・FSx要件の確認・合意", False, 11, None),
            ("2026年4月～          NW切替支援 着手（詳細設計・構築）", False, 11, None),
            ("2026年5月～          FSx for Windows File Server 設計開始", False, 11, None),
            ("2026年夏頃（目処）   NW切替完了", False, 11, None),
            ("2026年夏～秋         FSx for Windows 環境構築・動作確認", False, 11, None),
            ("2026年秋～冬         データ移行（DataSync）・権限確認", False, 11, None),
            ("2026年末～2027年初   カットオーバー・新環境運用開始", False, 11, None),
            ("2027年3月            既存機器 保守切れ", False, 11, None),
        ])
print("スケジュール更新")

# ===== STEP11: スライド順序を整える =====
# 現在(0-based): [0]表紙 [1]はじめに [2]前回VR [3]Win互換 [4]HA [5]FSxW単価
#                [6]費用サマリ [7]利用料概算 [8]今後section [9]確認事項
#                [10]スケジュール [11]裏表紙 [12]NW-header [13]NW-content
# 目標:
#  [0]表紙 [1]はじめに [2]前回VR [3]NW-header [4]NW-content
#  [5]Win互換 [6]HA [7]FSxW単価 [8]費用サマリ [9]利用料概算
#  [10]今後section [11]確認事項 [12]スケジュール [13]裏表紙
move_slide(prs, 12, 3)  # NW-header: 12→3
move_slide(prs, 13, 4)  # NW-content: 13→4

# 最終確認
print("\n最終スライド構成:")
for i, sl in enumerate(prs.slides):
    texts = [s.text_frame.text[:40].replace('\n',' ').replace('\x0b',' ')
             for s in sl.shapes
             if s.has_text_frame and s.text_frame.text.strip()
             and 'スライド番号' not in s.name]
    t = ' | '.join(texts[:2]) if texts else '(no text)'
    print(f"  [{i+1:02d}] {t}  [{sl.part.partname}]")

prs.save(dst)
print(f"\n保存完了: {dst}")
print(f"総スライド数: {len(prs.slides)}")
