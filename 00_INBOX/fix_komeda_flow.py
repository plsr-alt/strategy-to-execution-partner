from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from lxml import etree
import sys
sys.stdout.reconfigure(encoding='utf-8')

path = r"C:\Users\tshibasaki\Documents\03.案件情報\2025_案件\コメダ\01-プリセールス\012-提案資料\【コメダ様】ファイルサーバのクラウド移行提案資料_20260310_v2.pptx"
prs = Presentation(path)

# ===== ユーティリティ =====
def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    prs.slides.part.drop_rel(sldId.get(qn('r:id')))
    sldIdLst.remove(sldId)

def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    elem = sldIdLst[from_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(to_idx, elem)

def clear_para(para):
    p = para._p
    for tag in ['a:r', 'a:br', 'a:fld']:
        for e in p.findall(qn(tag)): p.remove(e)

def set_para_text(para, text):
    clear_para(para)
    if not text: return
    r = etree.SubElement(para._p, qn('a:r'))
    t = etree.SubElement(r, qn('a:t'))
    t.text = text

def add_para_with_text(tf, text, bold=False, size=None):
    para = tf.add_paragraph()
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    if size: run.font.size = Pt(size)
    return para

# ============================================================
# STEP1: 「今後の進め方について」セクションヘッダー(idx=10)を削除
# ============================================================
# 現構成確認
print("変更前:")
for i, sl in enumerate(prs.slides):
    txts = [s.text_frame.text[:30].replace(chr(11),' ') for s in sl.shapes if s.has_text_frame and s.text_frame.text.strip() and 'スライド番号' not in s.name]
    print(f"  [{i:02d}] {txts[0][:40] if txts else '(empty)'!r}")

target = "今後の進め方について"
delete_idx = None
for i, sl in enumerate(prs.slides):
    for s in sl.shapes:
        if s.has_text_frame and target in s.text_frame.text and i >= 9:
            # 最初のものだけ（セクションヘッダー）
            if delete_idx is None:
                delete_idx = i
                break

if delete_idx is not None:
    delete_slide(prs, delete_idx)
    print(f"\n[削除] idx={delete_idx} 「今後の進め方について」セクションヘッダー")

# 削除後の構成:
# [0]表紙 [1]はじめに [2]前回VR [3]NW-header [4]NW-content
# [5]Win互換 [6]HA [7]FSxW単価 [8]費用サマリ [9]利用料概算
# [10]確認事項 [11]スケジュール [12]裏表紙

# ============================================================
# STEP2: スライド順序を変更
# Win互換[5]→[3], HA[6]→[4], NW-header[3]→[5], NW-content[4]→[6]
# move(5,3): Win互換を3番へ
# ============================================================
# move Win互換(5) → 3
move_slide(prs, 5, 3)
# 状態: [0]表紙 [1]はじめに [2]前回VR [3]Win互換 [4]NW-header [5]NW-content [6]HA ...
# move HA(6) → 4
move_slide(prs, 6, 4)
# 状態: [0]表紙 [1]はじめに [2]前回VR [3]Win互換 [4]HA [5]NW-header [6]NW-content [7]FSxW単価 ...

print("\n並び替え完了:")
for i, sl in enumerate(prs.slides):
    txts = [s.text_frame.text[:35].replace(chr(11),' ') for s in sl.shapes if s.has_text_frame and s.text_frame.text.strip() and 'スライド番号' not in s.name]
    print(f"  [{i+1:02d}] {txts[0][:40] if txts else '(empty)'!r}")

# ============================================================
# STEP3: はじめに(idx=1) アジェンダを新しい流れに合わせて更新
# ============================================================
sl_intro = prs.slides[1]
for shape in sl_intro.shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'コンテンツ' in nm or 'Content' in nm or ('プレースホルダー' in nm and 'タイトル' not in nm and 'スライド' not in nm):
        tf = shape.text_frame
        # 既存テキストを確認して更新
        full_text = tf.text
        if 'アジェンダ' in full_text or 'NW切替' in full_text or '振り返り' in full_text:
            tf.clear()
            tf.word_wrap = True
            lines = [
                ("前回ご提案した内容を踏まえ、本日はFSx for Windows File Serverへの移行に向けた"
                 "要件確認・NW切替支援のご提案、および今後のスケジュール合意を目的としております。", False, 11),
                ("", False, 9),
                ("■ 本日のアジェンダ", True, 12),
                ("　１．前回ご提案内容の振り返り", False, 11),
                ("　２．FSx for Windows File Server 特徴・費用概算", False, 11),
                ("　３．NW切替支援について（移行の前提条件）", False, 11),
                ("　４．今後の進め方・スケジュール確認", False, 11),
                ("", False, 9),
                ("■ お打ち合わせのゴール", True, 12),
                ("　・FSxWの特徴と費用感をご確認いただく", False, 11),
                ("　・NW切替の進め方・スケジュールについて認識合わせする", False, 11),
                ("　・FSx設計開始（2026年5月〜）に向けた合意形成", False, 11),
            ]
            first = True
            for text, bold, size in lines:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                if text:
                    run = p.add_run()
                    run.text = text
                    run.font.bold = bold
                    run.font.size = Pt(size)
            print("はじめに アジェンダ更新")
            break

# ============================================================
# STEP4: NW-content(idx=6) 末尾に橋渡し文を追加
# ============================================================
sl_nw = prs.slides[6]
for shape in sl_nw.shapes:
    if not shape.has_text_frame: continue
    nm = shape.name
    if 'Content' in nm or 'コンテンツ' in nm or ('Placeholder' in nm and 'Title' not in nm):
        tf = shape.text_frame
        if 'NW切替が必要か' in tf.text or '接続方式' in tf.text:
            # 末尾に追記
            add_para_with_text(tf, "", False, 9)
            add_para_with_text(tf, "NW切替完了後のステップ", True, 13)
            add_para_with_text(tf, "  NW切替が完了次第、FSx for Windows File Serverの構築・移行フェーズへ移行します。", False, 11)
            add_para_with_text(tf, "  弊社がNW切替からFSx移行まで一貫してサポートいたします。", False, 11)
            print("NW-content 橋渡し文追加")
            break

# ============================================================
# 最終確認
# ============================================================
print("\n最終スライド構成:")
for i, sl in enumerate(prs.slides):
    txts = [s.text_frame.text[:40].replace(chr(11),' ')
            for s in sl.shapes
            if s.has_text_frame and s.text_frame.text.strip() and 'スライド番号' not in s.name]
    t = txts[0][:45] if txts else '(empty)'
    print(f"  [{i+1:02d}] {t!r}")

prs.save(path)
print(f"\n上書き保存完了: {path}")
print(f"総スライド数: {len(prs.slides)}")
