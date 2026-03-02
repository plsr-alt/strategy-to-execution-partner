#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
コメダ様 FSx選択肢メリデメ PPT生成スクリプト
python-pptx を使用して、FSx W/N/Qumuloの比較資料を自動生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

def add_title_slide(prs, title, subtitle):
    """タイトルスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # ブランクレイアウト

    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # 濃紺

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    # 日付
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"作成日: {datetime.now().strftime('%Y年%m月%d日')}"
    date_frame.paragraphs[0].font.size = Pt(14)
    date_frame.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)

def add_content_slide(prs, title, content_func):
    """コンテンツスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # ブランクレイアウト

    # タイトル背景（薄い青）
    title_bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(200, 220, 240)
    title_bg.line.color.rgb = RGBColor(0, 51, 102)

    # タイトルテキスト
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # コンテンツエリア
    content_func(slide)

def add_comparison_table(slide, data):
    """比較表を追加"""
    from pptx.util import Inches

    rows = len(data) + 1
    cols = len(data[0])

    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(9.4)
    height = Inches(4.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # ヘッダー
    headers = ["比較項目", "FSx Windows", "FSx NFS", "Qumulo"]
    for col_idx, header_text in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)

        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # データ行
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_text

            # 行替わりで背景色
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 255)

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.word_wrap = True

def add_text_box(slide, title, content_list, top_inches=1.2):
    """テキストボックス（箇条書き）を追加"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(top_inches), Inches(9), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for idx, item in enumerate(content_list):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(13)
        p.level = 0
        p.space_before = Pt(6)

def main():
    # プレゼンテーション作成
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ─────────────────────────────────────
    # スライド1: タイトル
    # ─────────────────────────────────────
    add_title_slide(prs, "AWS ファイルストレージ選択肢", "FSx Windows / NFS / Qumulo メリデメ比較")

    # ─────────────────────────────────────
    # スライド2: 全体比較表
    # ─────────────────────────────────────
    comparison_data = [
        ["費用（年間・50TB）", "約180万円", "約150万円", "別途見積"],
        ["運用難易度", "高⚠️", "中", "低✅"],
        ["学習曲線", "急", "中程度", "緩やか"],
        ["マネージド度", "完全マネージド", "完全マネージド", "マネージド（オプション）"],
        ["既存運用継承", "○", "○（部分的）", "○"],
        ["GUI操作", "○（Windows）", "○（System Manager）", "○（Web UI）"],
    ]

    def slide2_content(slide):
        add_comparison_table(slide, comparison_data)

    add_content_slide(prs, "3つの選択肢 — 全体比較", slide2_content)

    # ─────────────────────────────────────
    # スライド3: FSx for Windows File Server
    # ─────────────────────────────────────
    def slide3_content(slide):
        content = [
            "✅ メリット",
            "  • 既存のWindows運用方法をそのまま継承可能",
            "  • ユーザーが既存のWindowsファイル共有と同じ感覚で利用可能",
            "  • AD連携による細かい権限制御に対応",
            "",
            "⚠️  デメリット / 注意点",
            "  • 運用手順書（リストア・バックアップ）の作成がボリュームあり",
            "  • AWS側の管理・監視体制の構築が必要",
            "  • 費用が3選択肢の中で最も高い（約180万円/年）",
            "  • オンプレミアス側の人員スキルがそのままAWS側でも必要",
        ]
        add_text_box(slide, "", content)

    add_content_slide(prs, "選択肢1: FSx for Windows File Server", slide3_content)

    # ─────────────────────────────────────
    # スライド4: FSx for NetApp ONTAP
    # ─────────────────────────────────────
    def slide4_content(slide):
        content = [
            "✅ メリット",
            "  • NFS / SMB の両プロトコルをサポート（柔軟性が高い）",
            "  • NetApp ONTAP の豊富な機能（Snapshot、レプリケーション）",
            "  • 運用手順書は適度な複雑さ（FSx Windowsより簡潔）",
            "  • Snapshot活用でバックアップ・リストアが効率的",
            "  • 既存のNetApp環境がある場合、知見を活用可能",
            "",
            "⚠️  デメリット / 注意点",
            "  • ONTAP特有の管理・運用知識が必要（中程度の学習曲線）",
            "  • Windows運用とは異なるため、一部の既存手順は修正が必要",
            "  • 費用は中程度（約150万円/年）",
        ]
        add_text_box(slide, "", content)

    add_content_slide(prs, "選択肢2: FSx for NetApp ONTAP", slide4_content)

    # ─────────────────────────────────────
    # スライド5: Qumulo
    # ─────────────────────────────────────
    def slide5_content(slide):
        content = [
            "✅ メリット",
            "  • 運用の簡素化とコスト削減を同時に実現（ハイブリッド型）",
            "  • GUI ベースの直感的な管理画面",
            "  • Windows / NFS / SMB など複数プロトコルをサポート",
            "  • 既存運用の継承と新しい運用方法の折衷が可能",
            "",
            "⚠️  デメリット / 注意点",
            "  • AWS純正マネージドサービスではない（オンプレミアム型）",
            "  • 定期的なメンテナンス・パッチ適用が必要（Windowsより少ない）",
            "  • 初期費用・運用費用は別途見積が必要",
        ]
        add_text_box(slide, "", content)

    add_content_slide(prs, "選択肢3: Qumulo", slide5_content)

    # ─────────────────────────────────────
    # スライド6: 選択判定フロー
    # ─────────────────────────────────────
    def slide6_content(slide):
        content = [
            "【判定フロー】",
            "",
            "Q1: 既存のWindowsファイルサーバーの運用方法をそのまま引き継ぎたい？",
            "  → YES  →  FSx for Windows File Server 推奨",
            "  → NO   →  Q2へ",
            "",
            "Q2: NetApp ONTAP の豊富な機能（Snapshot等）を活用したい？",
            "     または、ONTAP環境がすでに存在する？",
            "  → YES  →  FSx for NetApp ONTAP 推奨",
            "  → NO   →  Q3へ",
            "",
            "Q3: 運用の簡素化と費用削減の両立を目指したい？",
            "     または、カスタマイズ性を重視したい？",
            "  → YES  →  Qumulo 推奨",
            "  → NO   →  要相談",
        ]
        add_text_box(slide, "", content)

    add_content_slide(prs, "選択肢の判定フロー", slide6_content)

    # ─────────────────────────────────────
    # スライド7: 次のステップ
    # ─────────────────────────────────────
    def slide7_content(slide):
        content = [
            "【次のステップ】",
            "",
            "① 本資料をご確認いただき、3つの選択肢からご検討ください。",
            "",
            "② ご不明な点、ご質問がございましたら、お気軽にお問い合わせください。",
            "",
            "③ お打ち合わせの際に以下を詳しくヒアリングさせていただきます：",
            "   • 移行後の運用方針（既存継承 vs 簡素化）",
            "   • 必要なサポート範囲（手順書、レクチャー、監視）",
            "   • AWS側の体制・リソース確保状況",
            "",
            "④ ヒアリング結果をもとに、詳細見積と提案書をお出しします。",
        ]
        add_text_box(slide, "", content)

    add_content_slide(prs, "次のステップ", slide7_content)

    # ファイル保存（WSLパス形式）
    import os
    import platform

    if platform.system() == "Linux":
        output_path = "/mnt/c/Users/tshibasaki/Desktop/etc/work/task/00_INBOX/コメダ様_FSx選択肢メリデメ_v1.pptx"
    else:
        output_path = r"C:\Users\tshibasaki\Desktop\etc\work\task\00_INBOX\コメダ様_FSx選択肢メリデメ_v1.pptx"

    output_dir = os.path.dirname(output_path)
    print(f"📁 OS: {platform.system()}")
    print(f"📁 出力ディレクトリ: {output_dir}")
    print(f"📁 ディレクトリ存在: {os.path.exists(output_dir)}")

    try:
        prs.save(output_path)
        saved_size = os.path.getsize(output_path)
        print(f"✅ PPT生成完了: {output_path}")
        print(f"✅ ファイルサイズ: {saved_size} bytes")
        print(f"✅ ファイル存在確認: {os.path.exists(output_path)}")
    except Exception as e:
        print(f"❌ エラー: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
